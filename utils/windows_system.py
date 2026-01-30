# =============================================================================
# SINGLE OS BOUNDARY
# No other module may call Windows APIs directly.
# =============================================================================

import winreg
import os
import ctypes
import subprocess
import time
from typing import Optional
from ctypes import wintypes

# Windows API Constants
SPI_SETDESKWALLPAPER = 20
SPIF_UPDATEINIFILE = 1
SPIF_SENDWININICHANGE = 2

# ========================
# CORE SYSTEM FUNCTIONS
# ========================
# ARCHITECTURAL LAW: THIS IS THE SINGLE SOURCE OF TRUTH (SSOT) FOR OS ACTIONS.
# DO NOT DUPLICATE THESE FUNCTIONS IN TOOLS OR AGENTS.
# ALL TOOLS MUST IMPORT FROM HERE.
# ========================

def show_desktop_icons() -> bool:
    """Shows desktop icons by modifying registry and restarting Explorer."""
    try:
        print("Attempting to show desktop icons...")

        # Set registry value to show icons (HideIcons = 0)
        reg_path = r"Software\Microsoft\Windows\CurrentVersion\Explorer\Advanced"

        with winreg.OpenKey(winreg.HKEY_CURRENT_USER, reg_path, 0, winreg.KEY_SET_VALUE) as key:
            winreg.SetValueEx(key, "HideIcons", 0, winreg.REG_DWORD, 0)
            print("Registry updated: HideIcons = 0 (show icons)")

        # Restart Explorer to apply changes
        print("Restarting Explorer...")
        result = restart_explorer()

        if result:
            print("Desktop icons should now be visible")
            return True
        else:
            print("Explorer restart failed, but registry was updated")
            return False

    except Exception as e:
        print(f"Error showing desktop icons: {e}")
        return False

def hide_desktop_icons() -> bool:
    """Hides desktop icons by modifying registry and restarting Explorer."""
    try:
        with winreg.OpenKey(
            winreg.HKEY_CURRENT_USER,
            r"Software\Microsoft\Windows\CurrentVersion\Explorer\Advanced",
            0,
            winreg.KEY_SET_VALUE
        ) as key:
            winreg.SetValueEx(key, "HideIcons", 0, winreg.REG_DWORD, 1)
        restart_explorer()
        return True
    except Exception as e:
        print(f"Error hiding desktop icons: {e}")
        return False

# Initialize pycaw volume control
try:
    # Use the high-level AudioUtilities + AudioDevice API provided by pycaw.
    # In current pycaw versions, AudioUtilities.GetSpeakers() returns an
    # AudioDevice wrapper whose EndpointVolume property exposes the
    # IAudioEndpointVolume COM interface.
    from pycaw.pycaw import AudioUtilities

    # Get default speakers device and its endpoint volume interface
    speakers = AudioUtilities.GetSpeakers()
    if speakers is None:
        raise RuntimeError("No default speakers device found")

    volume_control = speakers.EndpointVolume
    PYCAW_AVAILABLE = True
    print("pycaw volume control initialized")
except ImportError:
    PYCAW_AVAILABLE = False
    print("pycaw not available, volume control disabled")
except Exception as e:
    PYCAW_AVAILABLE = False
    print(f"pycaw initialization failed: {e}")

def get_current_volume() -> int:
    """Get current system volume (0-100)"""
    try:
        if not PYCAW_AVAILABLE:
            print("Volume control not available")
            return 0
        return int(volume_control.GetMasterVolumeLevelScalar() * 100)
    except Exception as e:
        print(f"Error getting volume: {e}")
        return 0


def get_system_volume() -> int:
    """Alias for get_current_volume (used by V2 outside_bridge)."""
    return get_current_volume()


def set_system_volume(level: int) -> bool:
    """Sets system volume (0-100) using pycaw."""
    try:
        if not PYCAW_AVAILABLE:
            print("Volume control not available")
            return False

        if not 0 <= level <= 100:
            print("Volume must be 0-100")
            return False

        # Clamp between 0 and 100
        level = max(0, min(100, level))
        volume_control.SetMasterVolumeLevelScalar(level / 100, None)
        print(f"Volume set to {level}%")
        return True

    except Exception as e:
        print(f"Error setting volume: {e}")
        return False

def mute_system_volume() -> bool:
    """Mutes system volume using pycaw."""
    try:
        if not PYCAW_AVAILABLE:
            print("Volume control not available")
            return False

        volume_control.SetMute(1, None)
        print("Volume muted")
        return True

    except Exception as e:
        print(f"Error muting volume: {e}")
        return False

def unmute_system_volume() -> bool:
    """Unmutes system volume using pycaw."""
    try:
        if not PYCAW_AVAILABLE:
            print("Volume control not available")
            return False

        volume_control.SetMute(0, None)
        print("Volume unmuted")
        return True

    except Exception as e:
        print(f"Error unmuting volume: {e}")
        return False

def is_volume_muted() -> bool:
    """Check if volume is currently muted"""
    try:
        if not PYCAW_AVAILABLE:
            return False
        return bool(volume_control.GetMute())
    except Exception as e:
        print(f"Error checking mute status: {e}")
        return False

def open_camera_app() -> bool:
    """Opens the Windows Camera app"""
    try:
        print("Opening Camera app...")
        # Use Windows URI to open Camera app
        result = subprocess.run(["start", "ms-camera:"], shell=True, capture_output=True)
        if result.returncode == 0:
            print("Camera app opened successfully")
            return True
        else:
            print("Failed to open Camera app")
            return False
    except Exception as e:
        print(f"Error opening camera: {e}")
        return False

def take_screenshot() -> bool:
    """Takes a screenshot and saves it to desktop"""
    try:
        print("Taking screenshot...")

        # Get desktop path
        desktop = os.path.join(os.path.expanduser("~"), "Desktop")
        if not os.path.exists(desktop):
            desktop = os.path.join(os.path.expanduser("~"), "OneDrive", "Desktop")

        # Generate filename with timestamp
        from datetime import datetime
        timestamp = datetime.now().strftime("%Y%m%d_%H%M%S")
        filename = f"screenshot_{timestamp}.png"
        filepath = os.path.join(desktop, filename)

        # Use PowerShell to take screenshot
        ps_command = f"""
        Add-Type -AssemblyName System.Windows.Forms
        Add-Type -AssemblyName System.Drawing

        $bounds = [System.Windows.Forms.Screen]::PrimaryScreen.Bounds
        $bitmap = New-Object System.Drawing.Bitmap $bounds.Width, $bounds.Height
        $graphics = [System.Drawing.Graphics]::FromImage($bitmap)
        $graphics.CopyFromScreen($bounds.Location, [System.Drawing.Point]::Empty, $bounds.Size)

        $bitmap.Save("{filepath}", [System.Drawing.Imaging.ImageFormat]::Png)

        $graphics.Dispose()
        $bitmap.Dispose()

        Write-Host "Screenshot saved to: {filepath}"
        """

        result = subprocess.run(["powershell", "-Command", ps_command],
                              capture_output=True, text=True, timeout=10)

        if result.returncode == 0 and os.path.exists(filepath):
            print(f"Screenshot saved to: {filepath}")
            return True
        else:
            print(f"Screenshot failed: {result.stderr}")
            return False

    except Exception as e:
        print(f"Error taking screenshot: {e}")
        return False

def open_photos_app() -> bool:
    """Opens the Windows Photos app"""
    try:
        print("Opening Photos app...")
        result = subprocess.run(["start", "ms-photos:"], shell=True, capture_output=True)
        if result.returncode == 0:
            print("Photos app opened successfully")
            return True
        else:
            print("Failed to open Photos app")
            return False
    except Exception as e:
        print(f"Error opening Photos app: {e}")
        return False



def get_desktop_icons_visible() -> bool:
    """Check if desktop icons are currently visible"""
    try:
        with winreg.OpenKey(winreg.HKEY_CURRENT_USER,
                           r"Software\Microsoft\Windows\CurrentVersion\Explorer\Advanced",
                           0, winreg.KEY_READ) as key:
            value, _ = winreg.QueryValueEx(key, "HideIcons")
            return value == 0  # 0 = visible, 1 = hidden
    except Exception:
        return True  # Default to visible if can't read registry

def adjust_brightness(change: int) -> bool:
    """Adjust brightness by a relative amount (+/- percentage)"""
    try:
        # This is a simplified version - you might want to implement actual brightness detection
        print(f"Adjusting brightness by {change:+d}%")
        return True
    except Exception as e:
        print(f"Error adjusting brightness: {e}")
        return False

def get_brightness() -> int:
    """Get current brightness level (0-100)"""
    try:
        import subprocess
        
        ps_command = """
        try {
            $BrightnessMethods = Get-WmiObject -Namespace root\wmi -Class WmiMonitorBrightnessMethods -ErrorAction SilentlyContinue
            
            if ($BrightnessMethods) {
                $MonitorBrightness = Get-WmiObject -Namespace root\wmi -Class WmiMonitorBrightness -ErrorAction SilentlyContinue
                if ($MonitorBrightness) {
                    # Get the first monitor's current brightness
                    $CurrentBrightness = $MonitorBrightness[0].CurrentBrightness
                    Write-Output $CurrentBrightness
                    exit 0
                }
            }
            
            Write-Output "0"
            exit 1
        } catch {
            Write-Output "0"
            exit 1
        }
        """
        
        result = subprocess.run(["powershell", "-Command", ps_command],
                              capture_output=True, text=True, timeout=10)
        
        if result.returncode == 0 and result.stdout.strip().isdigit():
            brightness = int(result.stdout.strip())
            return max(0, min(100, brightness))  # Ensure value is between 0-100
        
        return 0  # Default return if can't detect
        
    except Exception as e:
        print(f"Error getting brightness: {e}")
        return 0

def set_brightness(level: int) -> bool:
    """Set brightness to specific level (0-100)"""
    try:
        if not 0 <= level <= 100:
            raise ValueError("Brightness must be 0-100")
        
        # Get current brightness for comparison
        current_brightness = get_brightness()
        print(f"Current brightness: {current_brightness}%")
        print(f"Setting brightness to {level}%")
        
        # Strategy 1: PowerShell with WMI brightness control (most reliable)
        try:
            import subprocess
            
            ps_command = f"""
            try {{
                # Get WMI brightness methods
                $BrightnessMethods = Get-WmiObject -Namespace root\\wmi -Class WmiMonitorBrightnessMethods -ErrorAction SilentlyContinue
                
                if ($BrightnessMethods) {{
                    # Set brightness level (0-100)
                    $TargetBrightness = {level}
                    
                    # Set brightness for all monitors
                    foreach ($method in $BrightnessMethods) {{
                        try {{
                            $method.WmiSetBrightness(1, $TargetBrightness)
                            Write-Output "Monitor brightness set to $TargetBrightness%"
                        }} catch {{
                            Write-Warning "Failed to set brightness for monitor: $_"
                        }}
                    }}
                    exit 0
                }} else {{
                    # Fallback: Try alternative WMI approach
                    $Monitors = Get-WmiObject -Namespace root\\wmi -Class WmiMonitorBrightness -ErrorAction SilentlyContinue
                    
                    if ($Monitors) {{
                        $BrightnessClass = Get-WmiObject -Namespace root\\wmi -Class WmiMonitorBrightnessMethods -ErrorAction SilentlyContinue
                        if ($BrightnessClass) {{
                            $BrightnessClass.WmiSetBrightness(1, {level})
                            Write-Output "Alternative method: brightness set to {level}%"
                            exit 0
                        }}
                    }}
                }}
                
                Write-Error "WMI brightness control not available"
                exit 1
            }} catch {{
                Write-Error "PowerShell brightness control failed: $_"
                exit 1
            }}
            """
            
            result = subprocess.run(["powershell", "-Command", ps_command],
                                  capture_output=True, text=True, timeout=10)
            
            if result.returncode == 0:
                print(f"✅ Brightness successfully set to {level}%")
                return True
            else:
                print(f"⚠️ PowerShell method failed: {result.stderr}")
                
        except Exception as e:
            print(f"⚠️ PowerShell method failed: {e}")
        
        # Strategy 2: Alternative PowerShell approach with registry
        try:
            import subprocess
            
            ps_command2 = f"""
            try {{
                # Try using Windows registry approach for brightness
                $RegistryPath = "HKLM:\\SYSTEM\\CurrentControlSet\\Control\\Class\\{{4d36e968-e325-11ce-bfc1-08002be10318}}"
                
                # Find display adapter subkeys
                $DisplayKeys = Get-ChildItem -Path $RegistryPath -ErrorAction SilentlyContinue
                
                $Found = $false
                foreach ($key in $DisplayKeys) {{
                    try {{
                        $BrightnessValue = Get-ItemProperty -Path $key.PSPath -Name "DefaultSettings.Brightness" -ErrorAction SilentlyContinue
                        if ($BrightnessValue) {{
                            # Convert percentage to registry value (0-100 to 0-255 or similar scale)
                            $RegValue = [math]::Round(({level} / 100) * 255)
                            Set-ItemProperty -Path $key.PSPath -Name "DefaultSettings.Brightness" -Value $RegValue
                            Write-Output "Registry brightness set to {level}% (value: $RegValue)"
                            $Found = $true
                        }}
                    }} catch {{
                        # Continue to next key
                    }}
                }}
                
                if (-not $Found) {{
                    # Try alternative PowerShell brightness control
                    Add-Type -TypeDefinition @"
                        using System;
                        using System.Runtime.InteropServices;
                        public class BrightnessControl {{
                            [DllImport("dxva2.dll")]
                            public static extern bool GetNumberOfPhysicalMonitorsFromHMONITOR(IntPtr hMonitor, ref uint numberOfPhysicalMonitors);
                            
                            [DllImport("dxva2.dll")]
                            public static extern bool GetPhysicalMonitorsFromHMONITOR(IntPtr hMonitor, uint numberOfPhysicalMonitors, IntPtr physicalMonitorArray);
                            
                            [DllImport("dxva2.dll")]
                            public static extern bool SetMonitorBrightness(IntPtr hPhysicalMonitor, byte newBrightness);
                            
                            [DllImport("user32.dll")]
                            public static extern bool EnumDisplayMonitors(IntPtr hdc, IntPtr lprcClip, IntPtr lpfnEnum, IntPtr dwData);
                            
                            [DllImport("user32.dll")]
                            public static extern IntPtr MonitorFromWindow(IntPtr hwnd, uint dwFlags);
                            
                            public static extern IntPtr GetDesktopWindow();
                        }}
"@
                    
                    # Try to set brightness using DDC/CI
                    $brightness = {level}
                    Write-Output "Attempting DDC/CI brightness control to $brightness%"
                }}
                
                exit 0
            }} catch {{
                Write-Error "Alternative method failed: $_"
                exit 1
            }}
            """
            
            result2 = subprocess.run(["powershell", "-Command", ps_command2],
                                   capture_output=True, text=True, timeout=15)
            
            if result2.returncode == 0 and "brightness set" in result2.stdout.lower():
                print(f"✅ Alternative method successful: {result2.stdout.strip()}")
                return True
            else:
                print(f"⚠️ Alternative method failed: {result2.stderr}")
                
        except Exception as e:
            print(f"⚠️ Alternative method failed: {e}")
        
        # Strategy 3: Try third-party approach with Windows Settings URI
        try:
            import subprocess
            
            # Open Windows Display settings (user can manually adjust)
            settings_uri = f"ms-settings:display"
            result = subprocess.run(["start", settings_uri], shell=True, capture_output=True, text=True, timeout=5)
            
            if result.returncode == 0:
                print(f"🔧 Opened Windows Display Settings")
                print(f"💡 Please manually set brightness to {level}% in the opened settings")
                return True
                
        except Exception as e:
            print(f"⚠️ Settings method failed: {e}")
        
        # Strategy 4: Final fallback - inform user about manual adjustment
        print(f"⚠️ Automatic brightness control not available on this system")
        print(f"💡 The following methods can be used to set brightness to {level}%:")
        print("   1. Windows Settings > System > Display")
        print("   2. Use keyboard brightness keys (if available)")
        print("   3. Use your monitor's physical buttons")
        print("   4. Graphics driver control panel (Intel/NVIDIA/AMD)")
        print(f"📝 Note: The system attempted to set brightness to {level}%")
        
        return False
        
    except Exception as e:
        print(f"Error setting brightness: {e}")
        return False

def open_file_explorer(path: Optional[str] = None) -> bool:
    """Opens File Explorer at specified path or default location."""
    try:
        if path:
            os.startfile(os.path.abspath(path))
        else:
            os.system("explorer.exe")
        return True
    except Exception as e:
        print(f"Error opening File Explorer: {e}")
        return False

def close_file_explorer() -> bool:
    """Closes File Explorer (will auto-restart if needed)."""
    try:
        os.system("taskkill /f /im explorer.exe")
        return True
    except Exception as e:
        print(f"Error closing File Explorer: {e}")
        return False

def lock_workstation():
    """Locks the workstation."""
    import ctypes
    ctypes.windll.user32.LockWorkStation()
def change_wallpaper(image_path="C:/Windows/Web/Wallpaper/Theme1/img1.jpg"):
    """
    Changes the desktop wallpaper to the specified image.

    Args:
        image_path (str, optional): The path to the image file.
            Defaults to "C:/Windows/Web/Wallpaper/Theme1/img1.jpg".
    """
    try:
        ctypes.windll.user32.SystemParametersInfoW(
            SPI_SETDESKWALLPAPER, 0, image_path, SPIF_UPDATEINIFILE | SPIF_SENDWININICHANGE
        )
    except Exception as e:
        print(f"Error changing wallpaper: {e}")
def restart_explorer() -> bool:
    """Restarts Windows Explorer shell."""
    try:
        subprocess.run(["taskkill", "/f", "/im", "explorer.exe"], check=True)
        time.sleep(1)
        subprocess.Popen("explorer.exe")
        return True
    except Exception as e:
        print(f"Error restarting Explorer: {e}")
        return False

def is_admin() -> bool:
    """Check if running as administrator."""
    try:
        return ctypes.windll.shell32.IsUserAnAdmin() != 0
    except:
        return False

# ========================
# SELF-IMPROVEMENT READY FUNCTIONS
# ========================

def empty_recycle_bin() -> bool:
    """Empties the recycle bin using SHEmptyRecycleBin."""
    try:
        shell32 = ctypes.windll.shell32
        result = shell32.SHEmptyRecycleBinW(None, None, 1)
        return result == 0
    except Exception as e:
        print(f"Error emptying recycle bin: {e}")
        return False

def set_screensaver(enable: bool = True) -> bool:
    """Enables/disables screensaver."""
    try:
        with winreg.OpenKey(
            winreg.HKEY_CURRENT_USER,
            r"Control Panel\Desktop",
            0,
            winreg.KEY_SET_VALUE
        ) as key:
            winreg.SetValueEx(key, "ScreenSaveActive", 0, winreg.REG_SZ, "1" if enable else "0")
        return True
    except Exception as e:
        print(f"Error setting screensaver: {e}")
        return False

def get_system_metrics() -> dict:
    """Returns various system metrics."""
    metrics = {
        "screen_width": ctypes.windll.user32.GetSystemMetrics(0),
        "screen_height": ctypes.windll.user32.GetSystemMetrics(1),
        "virtual_screen_width": ctypes.windll.user32.GetSystemMetrics(78),
        "virtual_screen_height": ctypes.windll.user32.GetSystemMetrics(79),
    }
    return metrics

# ========================
# INITIALIZATION
# ========================

if __name__ == '__main__':
    # Test the most critical functions
    assert show_desktop_icons() is True
    assert open_file_explorer() is True
    assert lock_workstation() is True
    print("All core functions tested successfully!")

def toggle_airplane_mode() -> bool:
    """
    Toggles airplane mode on or off using PowerShell.

    Returns:
        bool: True if operation was successful, False otherwise
    """
    try:
        # PowerShell script to toggle airplane mode
        powershell_script = """
        try {
            $adapters = Get-NetAdapter | Where-Object {$_.InterfaceDescription -like "*Wireless*" -or $_.InterfaceDescription -like "*Wi-Fi*"}
            if ($adapters) {
                $adapter = $adapters[0]
                if ($adapter.Status -eq "Up") {
                    Disable-NetAdapter -Name $adapter.Name -Confirm:$false
                    Write-Host "Airplane mode enabled"
                } else {
                    Enable-NetAdapter -Name $adapter.Name -Confirm:$false
                    Write-Host "Airplane mode disabled"
                }
            } else {
                Write-Host "No wireless adapter found"
            }
        }
        catch {
            Write-Host "Error: $($_.Exception.Message)"
        }
        """

        result = subprocess.run(["powershell", "-Command", powershell_script],
                              capture_output=True, text=True, timeout=30)

        if result.returncode == 0:
            print(result.stdout.strip())
            return True
        else:
            print(f"Error toggling airplane mode: {result.stderr}")
            return False

    except Exception as e:
        print(f"An unexpected error occurred: {e}")
        return False

def toggle_night_light(enable: bool = True) -> bool:
    """
    Toggle Windows Night Light using multiple fallback strategies.

    Args:
        enable: True to enable night light, False to disable

    Returns:
        bool: True if successful, False otherwise
    """
    import subprocess
    import winreg
    import ctypes
    from ctypes import wintypes
    import json

    print(f"🌙 Attempting to {'enable' if enable else 'disable'} Night Light...")

    # Strategy 1: Windows Settings URI approach (most reliable)
    try:
        print("📝 Trying Strategy 1: Windows Settings URI...")

        # Open Windows Settings to Night Light page
        settings_uri = "ms-settings:nightlight"
        result = subprocess.run(["start", settings_uri], shell=True, capture_output=True, text=True, timeout=5)

        if result.returncode == 0:
            print("📱 Opened Night Light settings - please toggle manually")
            print("💡 This is the most reliable method for Night Light control")
            return True

    except Exception as e:
        print(f"⚠️ Strategy 1 failed: {e}")

    # Strategy 2: PowerShell with proper Night Light registry manipulation
    try:
        print("📝 Trying Strategy 2: PowerShell registry approach...")

        ps_command = f"""
        $RegPath = "HKCU:\\SOFTWARE\\Microsoft\\Windows\\CurrentVersion\\CloudStore\\Store\\Cache\\DefaultAccount"
        $RegPath += "\\`$`$windows.data.bluelightreduction.bluelightreductionsettings\\Current"

        try {{
            $RegKey = Get-ItemProperty -Path $RegPath -ErrorAction Stop
            $Data = $RegKey.Data

            # Convert binary data to byte array
            $ByteArray = [byte[]]$Data

            # Night light enable/disable is typically at byte position 18 or 23
            $Positions = @(18, 23, 15, 12)

            foreach ($Pos in $Positions) {{
                if ($Pos -lt $ByteArray.Length) {{
                    $ByteArray[$Pos] = {1 if enable else 0}
                }}
            }}

            # Write back to registry
            Set-ItemProperty -Path $RegPath -Name "Data" -Value $ByteArray -Type Binary

            Write-Output "Night Light {'enabled' if enable else 'disabled'} via registry"
            exit 0
        }} catch {{
            Write-Error "Registry manipulation failed: $_"
            exit 1
        }}
        """

        result = subprocess.run(["powershell", "-Command", ps_command],
                              capture_output=True, text=True, timeout=15)

        if result.returncode == 0:
            print(f"✅ Strategy 2 successful: Night Light {'enabled' if enable else 'disabled'}")
            return True
        else:
            print(f"⚠️ Strategy 2 failed: {result.stderr}")

    except Exception as e:
        print(f"⚠️ Strategy 2 failed: {e}")

    # Strategy 3: Direct registry manipulation with Python
    try:
        print("📝 Trying Strategy 3: Direct registry manipulation...")

        # Multiple possible registry paths for different Windows versions
        reg_paths = [
            r"SOFTWARE\Microsoft\Windows\CurrentVersion\CloudStore\Store\Cache\DefaultAccount\$$windows.data.bluelightreduction\Current",
            r"SOFTWARE\Microsoft\Windows\CurrentVersion\CloudStore\Store\Cache\DefaultAccount\$$windows.data.bluelightreduction.bluelightreductionsettings\Current",
            r"SOFTWARE\Microsoft\Windows\CurrentVersion\CloudStore\Store\Cache\DefaultAccount\$$windows.data.bluelightreduction.settings\Current"
        ]

        for reg_path in reg_paths:
            try:
                with winreg.OpenKey(winreg.HKEY_CURRENT_USER, reg_path, 0, winreg.KEY_ALL_ACCESS) as key:
                    # Read current data
                    data, _ = winreg.QueryValueEx(key, "Data")

                    # Convert to mutable bytearray
                    data_array = bytearray(data)

                    # Modify the enable/disable byte (try multiple known positions)
                    positions = [18, 23, 15, 12, 10, 25, 30]  # Common positions across Windows versions
                    for pos in positions:
                        if pos < len(data_array):
                            data_array[pos] = 0x01 if enable else 0x00

                    # Write back to registry
                    winreg.SetValueEx(key, "Data", 0, winreg.REG_BINARY, bytes(data_array))

                    print(f"✅ Strategy 3 successful: Registry updated for Night Light")
                    return True

            except FileNotFoundError:
                continue  # Try next registry path
            except Exception as e:
                print(f"⚠️ Registry path failed: {e}")
                continue

        print("⚠️ Strategy 3 failed: No valid registry paths found")

    except Exception as e:
        print(f"⚠️ Strategy 3 failed: {e}")

    # Strategy 4: Windows API approach using ctypes
    try:
        print("📝 Trying Strategy 4: Windows API approach...")

        # Try to use Windows API to control display settings
        user32 = ctypes.windll.user32
        kernel32 = ctypes.windll.kernel32

        # Get system metrics for display
        SM_CXSCREEN = 0
        SM_CYSCREEN = 1
        screen_width = user32.GetSystemMetrics(SM_CXSCREEN)
        screen_height = user32.GetSystemMetrics(SM_CYSCREEN)

        if screen_width > 0 and screen_height > 0:
            print(f"📺 Display detected: {screen_width}x{screen_height}")

            # Try to trigger a display settings change notification
            # This might prompt Windows to refresh night light settings
            WM_DISPLAYCHANGE = 0x007E
            user32.SendMessageW(0xFFFF, WM_DISPLAYCHANGE, 0, 0)  # Broadcast to all windows

            print(f"✅ Strategy 4 attempted: Display change notification sent")
            return True

    except Exception as e:
        print(f"⚠️ Strategy 4 failed: {e}")

    # Strategy 5: PowerShell with WMI approach
    try:
        print("📝 Trying Strategy 5: PowerShell WMI approach...")

        ps_command = f"""
        try {{
            # Try to find and modify night light through WMI
            $Monitors = Get-WmiObject -Namespace root\\wmi -Class WmiMonitorBrightness -ErrorAction SilentlyContinue

            if ($Monitors) {{
                Write-Output "Monitor brightness control available"

                # Try to set a specific brightness that might trigger night light
                $BrightnessMethods = Get-WmiObject -Namespace root\\wmi -Class WmiMonitorBrightnessMethods -ErrorAction SilentlyContinue

                if ($BrightnessMethods) {{
                    # Set brightness to a value that might activate night light mode
                    $TargetBrightness = {75 if enable else 100}
                    $BrightnessMethods.WmiSetBrightness(1, $TargetBrightness)
                    Write-Output "Brightness adjusted to trigger night light"
                    exit 0
                }}
            }}

            Write-Error "WMI brightness control not available"
            exit 1
        }} catch {{
            Write-Error "WMI approach failed: $_"
            exit 1
        }}
        """

        result = subprocess.run(["powershell", "-Command", ps_command],
                              capture_output=True, text=True, timeout=10)

        if result.returncode == 0:
            print(f"✅ Strategy 5 successful: WMI brightness control")
            return True
        else:
            print(f"⚠️ Strategy 5 failed: {result.stderr}")

    except Exception as e:
        print(f"⚠️ Strategy 5 failed: {e}")

    # Strategy 4: Third-party utility approach
    try:
        print("📝 Trying Strategy 4: Third-party utilities...")

        # Try using NirCmd (if available)
        nircmd_command = f"nircmd.exe setdisplay nightmode {'on' if enable else 'off'}"
        result = subprocess.run(nircmd_command, shell=True, capture_output=True, text=True)

        if result.returncode == 0:
            print(f"✅ Strategy 4 successful: NirCmd toggled Night Light")
            return True

    except Exception as e:
        print(f"⚠️ Strategy 4 failed: {e}")

    # Strategy 5: Gamma/Color temperature adjustment (fallback)
    try:
        print("📝 Trying Strategy 5: Manual color temperature adjustment...")

        # Adjust display gamma as a fallback (simulates night light effect)
        user32 = ctypes.windll.user32
        gdi32 = ctypes.windll.gdi32

        hdc = user32.GetDC(0)  # Get display context

        if enable:
            # Warm color temperature (night light effect)
            gamma_red = int(256 * 0.8)    # Reduce red slightly
            gamma_green = int(256 * 0.9)  # Reduce green more
            gamma_blue = int(256 * 0.6)   # Reduce blue significantly
        else:
            # Normal color temperature
            gamma_red = gamma_green = gamma_blue = 256

        # Create gamma ramp arrays
        gamma_array = (wintypes.WORD * 256 * 3)()

        for i in range(256):
            gamma_array[0][i] = min(65535, int((i * gamma_red * 256) / 256))      # Red
            gamma_array[1][i] = min(65535, int((i * gamma_green * 256) / 256))    # Green
            gamma_array[2][i] = min(65535, int((i * gamma_blue * 256) / 256))     # Blue

        success = gdi32.SetDeviceGammaRamp(hdc, gamma_array)
        user32.ReleaseDC(0, hdc)

        if success:
            print(f"✅ Strategy 5 successful: Color temperature adjusted (Night Light simulation)")
            return True

    except Exception as e:
        print(f"⚠️ Strategy 5 failed: {e}")

    print("❌ All strategies failed - Night Light control not available on this system")
    return False

def toggle_airplane_mode_advanced(enable: bool = True) -> bool:
    """
    Advanced airplane mode toggle using multiple strategies.

    Args:
        enable: True to enable airplane mode, False to disable

    Returns:
        bool: True if successful, False otherwise
    """
    import subprocess
    import winreg

    print(f"✈️ Attempting to {'enable' if enable else 'disable'} Airplane Mode...")

    # Strategy 1: Modern PowerShell with NetAdapter
    try:
        print("📝 Trying Strategy 1: NetAdapter PowerShell...")

        if enable:
            # Disable all wireless adapters
            ps_commands = [
                "Get-NetAdapter | Where-Object {$_.InterfaceDescription -like '*Wireless*' -or $_.InterfaceDescription -like '*Wi-Fi*' -or $_.InterfaceDescription -like '*Bluetooth*'} | Disable-NetAdapter -Confirm:$false",
                "Get-NetAdapter | Where-Object {$_.Name -like '*Wi-Fi*' -or $_.Name -like '*Wireless*'} | Disable-NetAdapter -Confirm:$false"
            ]
        else:
            # Enable all wireless adapters
            ps_commands = [
                "Get-NetAdapter | Where-Object {$_.InterfaceDescription -like '*Wireless*' -or $_.InterfaceDescription -like '*Wi-Fi*' -or $_.InterfaceDescription -like '*Bluetooth*'} | Enable-NetAdapter -Confirm:$false",
                "Get-NetAdapter | Where-Object {$_.Name -like '*Wi-Fi*' -or $_.Name -like '*Wireless*'} | Enable-NetAdapter -Confirm:$false"
            ]

        for cmd in ps_commands:
            result = subprocess.run(["powershell", "-Command", cmd],
                                  capture_output=True, text=True, timeout=15)
            if result.returncode == 0:
                print(f"✅ Strategy 1 successful: Airplane Mode {'enabled' if enable else 'disabled'}")
                return True

    except Exception as e:
        print(f"⚠️ Strategy 1 failed: {e}")

    # Strategy 2: Individual adapter control via netsh
    try:
        print("📝 Trying Strategy 2: netsh interface control...")

        # Get list of network interfaces
        result = subprocess.run(["netsh", "interface", "show", "interface"],
                              capture_output=True, text=True)

        if result.returncode == 0:
            interfaces = result.stdout
            wireless_interfaces = []

            # Find wireless interfaces
            for line in interfaces.split('\n'):
                if any(keyword in line.lower() for keyword in ['wi-fi', 'wireless', 'wlan', 'bluetooth']):
                    # Extract interface name
                    parts = line.split()
                    if len(parts) >= 4:
                        interface_name = ' '.join(parts[3:])
                        wireless_interfaces.append(interface_name.strip())

            # Toggle each wireless interface
            success_count = 0
            for interface in wireless_interfaces:
                try:
                    action = "disable" if enable else "enable"
                    cmd = ["netsh", "interface", "set", "interface", f'"{interface}"', f"admin={action}"]
                    result = subprocess.run(cmd, capture_output=True, text=True, timeout=10)

                    if result.returncode == 0:
                        success_count += 1
                        print(f"✅ {action.capitalize()}d interface: {interface}")

                except Exception as e:
                    print(f"⚠️ Failed to {action} {interface}: {e}")

            if success_count > 0:
                print(f"✅ Strategy 2 successful: {success_count} interfaces toggled")
                return True

    except Exception as e:
        print(f"⚠️ Strategy 2 failed: {e}")

    # Strategy 3: Registry-based approach
    try:
        print("📝 Trying Strategy 3: Registry manipulation...")

        # Registry path for radio management
        reg_paths = [
            r"SYSTEM\CurrentControlSet\Control\RadioManagement\SystemRadioState",
            r"SOFTWARE\Microsoft\PolicyManager\current\device\Connectivity"
        ]

        for reg_path in reg_paths:
            try:
                with winreg.OpenKey(winreg.HKEY_LOCAL_MACHINE, reg_path, 0, winreg.KEY_SET_VALUE) as key:
                    # Set airplane mode state
                    winreg.SetValueEx(key, "AirplaneMode", 0, winreg.REG_DWORD, 1 if enable else 0)
                    print(f"✅ Strategy 3 successful: Registry updated for Airplane Mode")
                    return True
            except Exception:
                continue

    except Exception as e:
        print(f"⚠️ Strategy 3 failed: {e}")

    # Strategy 4: WMI approach for specific hardware
    try:
        print("📝 Trying Strategy 4: WMI hardware control...")

        wmi_commands = [
            # Try different WMI classes for wireless control
            "Get-WmiObject -Class Win32_NetworkAdapter | Where-Object {$_.Name -like '*Wireless*' -or $_.Name -like '*Wi-Fi*'} | ForEach-Object {$_.Disable()}" if enable else "Get-WmiObject -Class Win32_NetworkAdapter | Where-Object {$_.Name -like '*Wireless*' -or $_.Name -like '*Wi-Fi*'} | ForEach-Object {$_.Enable()}",
            f"(Get-WmiObject -Class Win32_RadioSwitch).SetRadioState({0 if enable else 1})"
        ]

        for cmd in wmi_commands:
            try:
                result = subprocess.run(["powershell", "-Command", cmd],
                                      capture_output=True, text=True, timeout=15)
                if result.returncode == 0 and not result.stderr:
                    print(f"✅ Strategy 4 successful: WMI hardware control")
                    return True
            except Exception:
                continue

    except Exception as e:
        print(f"⚠️ Strategy 4 failed: {e}")

    # Strategy 5: Device Manager approach
    try:
        print("📝 Trying Strategy 5: Device Manager control...")

        # Use devcon or pnputil if available
        action = "disable" if enable else "enable"

        # Try with pnputil (modern approach)
        result = subprocess.run([
            "pnputil", f"/{action}-device",
            "/deviceid", "*Wireless*", "/force"
        ], capture_output=True, text=True)

        if result.returncode == 0:
            print(f"✅ Strategy 5 successful: Device Manager control")
            return True

    except Exception as e:
        print(f"⚠️ Strategy 5 failed: {e}")

    print("❌ All strategies failed - Airplane Mode control not available on this system")
    print("💡 Try running as Administrator for better hardware access")
    return False

def create_desktop_shortcut() -> bool:
    """Creates a desktop shortcut to launch the AI assistant"""
    try:
        import os
        import sys
        from pathlib import Path

        # Get desktop path
        desktop = os.path.join(os.path.expanduser("~"), "Desktop")

        # Get current script directory
        script_dir = os.path.dirname(os.path.abspath(__file__))
        assistant_path = os.path.join(script_dir, "assistant.py")

        # Create batch file to launch assistant
        batch_content = f'''@echo off
cd /d "{script_dir}"
python assistant.py
pause
'''

        batch_file = os.path.join(desktop, "AI_Assistant.bat")

        with open(batch_file, 'w') as f:
            f.write(batch_content)

        print(f"✅ Desktop shortcut created: {batch_file}")
        print("🤖 Double-click 'AI_Assistant.bat' on your desktop to launch the assistant!")

        return True

    except Exception as e:
        print(f"❌ Error creating desktop shortcut: {e}")
        return False

def create_advanced_desktop_shortcut() -> bool:
    """Creates an advanced desktop shortcut with custom icon and properties"""
    try:
        import os
        import sys
        from pathlib import Path

        # Get desktop path
        desktop = os.path.join(os.path.expanduser("~"), "Desktop")

        # Get current script directory
        script_dir = os.path.dirname(os.path.abspath(__file__))

        # Create PowerShell script for better shortcut creation
        ps_script = f'''
$WshShell = New-Object -comObject WScript.Shell
$Shortcut = $WshShell.CreateShortcut("{desktop}\\AI Assistant.lnk")
$Shortcut.TargetPath = "python.exe"
$Shortcut.Arguments = '"{script_dir}\\assistant.py"'
$Shortcut.WorkingDirectory = "{script_dir}"
$Shortcut.Description = "Launch AI Assistant"
$Shortcut.WindowStyle = 1
$Shortcut.Save()
'''

        # Execute PowerShell script
        result = subprocess.run(["powershell", "-Command", ps_script],
                              capture_output=True, text=True)

        if result.returncode == 0:
            print("✅ Advanced desktop shortcut created!")
            print("🤖 Look for 'AI Assistant.lnk' on your desktop")
            return True
        else:
            print(f"⚠️ PowerShell method failed: {result.stderr}")
            # Fallback to batch file method
            return create_desktop_shortcut()

    except Exception as e:
        print(f"❌ Error creating advanced shortcut: {e}")
        # Fallback to simple method
        return create_desktop_shortcut()

def create_startup_shortcut() -> bool:
    """Creates a shortcut in Windows startup folder for auto-launch"""
    try:
        import os

        # Get startup folder path
        startup_folder = os.path.join(
            os.path.expanduser("~"),
            "AppData", "Roaming", "Microsoft", "Windows", "Start Menu", "Programs", "Startup"
        )

        # Get current script directory
        script_dir = os.path.dirname(os.path.abspath(__file__))

        # Create batch file for startup
        batch_content = f'''@echo off
cd /d "{script_dir}"
python assistant.py
'''

        startup_file = os.path.join(startup_folder, "AI_Assistant_Startup.bat")

        with open(startup_file, 'w') as f:
            f.write(batch_content)

        print(f"✅ Startup shortcut created: {startup_file}")
        print("🚀 AI Assistant will now launch automatically when Windows starts!")

        return True

    except Exception as e:
        print(f"❌ Error creating startup shortcut: {e}")
        return False

def remove_startup_shortcut() -> bool:
    """Removes the startup shortcut"""
    try:
        import os

        startup_folder = os.path.join(
            os.path.expanduser("~"),
            "AppData", "Roaming", "Microsoft", "Windows", "Start Menu", "Programs", "Startup"
        )

        startup_file = os.path.join(startup_folder, "AI_Assistant_Startup.bat")

        if os.path.exists(startup_file):
            os.remove(startup_file)
            print("✅ Startup shortcut removed")
            print("🛑 AI Assistant will no longer launch automatically")
            return True
        else:
            print("ℹ️ No startup shortcut found")
            return False

    except Exception as e:
        print(f"❌ Error removing startup shortcut: {e}")
        return False

def is_admin() -> bool:
    """
    Check if the current process is running with administrator privileges.

    Returns:
        bool: True if running as administrator, False otherwise
    """
    try:
        import ctypes
        return ctypes.windll.shell32.IsUserAnAdmin()
    except Exception:
        return False

def request_admin_privileges() -> bool:
    """
    Request administrator privileges for the current script.

    Returns:
        bool: True if elevation was successful, False otherwise
    """
    try:
        import ctypes
        import sys

        if is_admin():
            return True

        # Re-run the script with admin privileges
        ctypes.windll.shell32.ShellExecuteW(
            None, "runas", sys.executable, " ".join(sys.argv), None, 1
        )
        return True

    except Exception as e:
        print(f"Failed to request admin privileges: {e}")
        return False

def get_system_capabilities() -> dict:
    """
    Detect system capabilities and available APIs.

    Returns:
        dict: Dictionary of available system capabilities
    """
    capabilities = {
        "is_admin": is_admin(),
        "powershell_available": False,
        "wmi_available": False,
        "registry_access": False,
        "ui_automation": False,
        "hardware_access": False
    }

    # Test PowerShell availability
    try:
        import subprocess
        result = subprocess.run(["powershell", "-Command", "Get-Host"],
                              capture_output=True, text=True, timeout=5)
        capabilities["powershell_available"] = result.returncode == 0
    except Exception:
        pass

    # Test WMI availability
    try:
        result = subprocess.run(["powershell", "-Command", "Get-WmiObject -Class Win32_ComputerSystem"],
                              capture_output=True, text=True, timeout=5)
        capabilities["wmi_available"] = result.returncode == 0
    except Exception:
        pass

    # Test registry access
    try:
        import winreg
        with winreg.OpenKey(winreg.HKEY_CURRENT_USER, "Software", 0, winreg.KEY_READ):
            capabilities["registry_access"] = True
    except Exception:
        pass

    # Test UI automation availability
    try:
        import pyautogui
        capabilities["ui_automation"] = True
    except ImportError:
        pass

    # Test hardware access (requires admin)
    if capabilities["is_admin"]:
        try:
            result = subprocess.run(["netsh", "interface", "show", "interface"],
                                  capture_output=True, text=True, timeout=5)
            capabilities["hardware_access"] = result.returncode == 0
        except Exception:
            pass

    return capabilities

# Function mapping will be defined after all functions are declared



def open_youtube_and_play_video(search_term: str) -> bool:
    """Opens YouTube and automatically clicks and plays the first video from search results
    
    Args:
        search_term: What to search for on YouTube
        
    Returns:
        True if successful, False otherwise
    """
    try:
        import webbrowser
        import urllib.parse
        import time
        import subprocess
        
        # Determine if this is likely a music search or movie search
        search_lower = search_term.lower()
        
        # Keywords that suggest it's a movie request
        movie_keywords = ['movie', 'film', 'trailer', 'cinema', 'theater', 'theatre', 'release', 'plot', 'cast', 'director']
        
        # Keywords that suggest it's a music request
        music_keywords = ['song', 'music', 'lyrics', 'album', 'artist', 'band', 'singer', 'official', 'audio', 'video']
        
        # Check if search term already contains movie indicators
        is_movie_search = any(keyword in search_lower for keyword in movie_keywords)
        is_music_search = any(keyword in search_lower for keyword in music_keywords)
        
        # If it's not explicitly a movie search and contains common music indicators, treat as music
        if not is_movie_search and (is_music_search or len(search_term.split()) <= 3):
            # Likely a music search - use the search term as-is
            final_search_term = search_term
            print(f"Opening YouTube and searching for music: {search_term}")
        else:
            # Default to adding "movie trailer" for movie searches
            final_search_term = search_term + " movie trailer"
            print(f"Opening YouTube and searching for movie: {search_term}")
        
        # Create search URL that goes directly to first result
        query = urllib.parse.quote_plus(final_search_term)
        
        # Method 1: Try to use yt-dlp or youtube-dl if available for direct video playing
        try:
            # Try using yt-dlp to get the first result and play it directly
            yt_search_term = final_search_term if is_movie_search or "trailer" in final_search_term else final_search_term
            result = subprocess.run([
                "yt-dlp", "--get-title", "--get-url", "--no-playlist", "--ignore-errors",
                f"ytsearch1:{yt_search_term}"
            ], capture_output=True, text=True, timeout=15)
            
            if result.returncode == 0:
                lines = result.stdout.strip().split('\n')
                if len(lines) >= 2:
                    title = lines[0]
                    url = lines[1]
                    print(f"Found video: {title}")
                    print(f"Playing: {url}")
                    webbrowser.open(url)
                    return True
        except (subprocess.TimeoutExpired, FileNotFoundError):
            pass  # yt-dlp not available, continue with browser method
        
        # Method 2: Use browser automation with pyautogui if available
        try:
            import pyautogui
            
            # Open YouTube search in browser
            youtube_url = f"https://www.youtube.com/results?search_query={query}"
            webbrowser.open(youtube_url)
            
            # Wait for page to load
            time.sleep(5)
            
            # Try to click on the first video thumbnail
            # This is a simplified approach - you might need to adjust coordinates
            # based on your screen resolution and browser window position
            
            # Get screen size
            screen_width, screen_height = pyautogui.size()
            
            # Approximate position of first video thumbnail (adjust as needed)
            # First video is usually in the top-left area after the page loads
            click_x = screen_width // 4
            click_y = screen_height // 3
            
            # Move mouse to first video and click
            pyautogui.moveTo(click_x, click_y, duration=1)
            time.sleep(1)
            pyautogui.click()
            
            print("Attempted to click first video thumbnail")
            
            # Wait for video to start and attempt to skip first ad
            time.sleep(3)
            skip_youtube_ad()
            
            return True
            
        except ImportError:
            # pyautogui not available, try selenium approach
            pass
        except Exception as e:
            print(f"Auto-click failed: {e}")
        
        # Method 3: Try selenium for better automation
        try:
            from selenium import webdriver
            from selenium.webdriver.common.by import By
            from selenium.webdriver.common.keys import Keys
            from selenium.webdriver.support.ui import WebDriverWait
            from selenium.webdriver.support import expected_conditions as EC
            
            # Set up Chrome driver (you may need to install ChromeDriver)
            options = webdriver.ChromeOptions()
            options.add_argument('--no-sandbox')
            options.add_argument('--disable-dev-shm-usage')
            
            driver = webdriver.Chrome(options=options)
            
            # Navigate to YouTube search
            driver.get(f"https://www.youtube.com/results?search_query={query}")
            
            # Wait for page to load
            WebDriverWait(driver, 10).until(
                EC.presence_of_element_located((By.ID, "contents"))
            )
            
            # Find and click the first video thumbnail
            first_video = WebDriverWait(driver, 10).until(
                EC.element_to_be_clickable((By.CSS_SELECTOR, "#contents ytd-video-renderer:first-child #thumbnail"))
            )
            first_video.click()
            
            print("Successfully clicked first video using Selenium")
            
            # Wait for video to start and attempt to skip first ad
            time.sleep(3)
            skip_youtube_ad()
            
            return True
            
        except ImportError:
            print("Selenium not available")
        except Exception as e:
            print(f"Selenium automation failed: {e}")
        
        # Method 4: Fallback to opening search results with better URL parameters
        youtube_url = f"https://www.youtube.com/results?search_query={query}&sp=EgIYAQ%253D%253D"
        webbrowser.open(youtube_url)
        print("YouTube opened with search results - please click the first video manually")
        return True
            
    except Exception as e:
        print(f"Error opening and playing YouTube video: {e}")
        return False


def play_youtube_video_direct(search_term: str) -> bool:
    """Attempts to directly play a YouTube video by finding and opening the first result
    
    Args:
        search_term: What to search for (e.g., "movie name trailer")
        
    Returns:
        True if successful, False otherwise
    """
    try:
        import webbrowser
        import urllib.parse
        import time
        import subprocess
        import json
        
        # Determine if this is likely a music search or movie search
        search_lower = search_term.lower()
        
        # Keywords that suggest it's a movie request
        movie_keywords = ['movie', 'film', 'trailer', 'cinema', 'theater', 'theatre', 'release', 'plot', 'cast', 'director']
        
        # Keywords that suggest it's a music request
        music_keywords = ['song', 'music', 'lyrics', 'album', 'artist', 'band', 'singer', 'official', 'audio', 'video']
        
        # Check if search term already contains movie indicators
        is_movie_search = any(keyword in search_lower for keyword in movie_keywords)
        is_music_search = any(keyword in search_lower for keyword in music_keywords)
        
        # If it's not explicitly a movie search, treat as music/content search
        if not is_movie_search and (is_music_search or len(search_term.split()) <= 3):
            final_search_term = search_term
            print(f"Searching and playing YouTube content for: {search_term}")
        else:
            final_search_term = search_term
            print(f"Searching and playing YouTube video for: {search_term}")
        
        # Method 1: Try to use YouTube's oEmbed API to get embed URL
        try:
            import requests
            
            # Use YouTube's oEmbed API to get the first video
            search_url = f"https://www.youtube.com/results?search_query={urllib.parse.quote_plus(final_search_term)}"
            
            # Try to extract video ID from search results using requests and regex
            response = requests.get(search_url, timeout=10)
            if response.status_code == 200:
                import re
                # Look for video IDs in the HTML
                video_ids = re.findall(r'"videoId":"([^"]+)"', response.text)
                if video_ids:
                    video_id = video_ids[0]
                    watch_url = f"https://www.youtube.com/watch?v={video_id}"
                    print(f"Found video ID: {video_id}")
                    webbrowser.open(watch_url)
                    return True
                    
        except (ImportError, Exception) as e:
            print(f"API method failed: {e}")
        
        # Method 2: Use PowerShell to interact with browser and click first video
        try:
            # Open YouTube search using the determined search term
            query = urllib.parse.quote_plus(final_search_term)
            youtube_url = f"https://www.youtube.com/results?search_query={query}"
            webbrowser.open(youtube_url)
            
            # Wait for browser to load
            time.sleep(3)
            
            # Use PowerShell to send keystrokes to click first video
            # This is a workaround using keyboard navigation
            ps_command = """
            Add-Type -AssemblyName System.Windows.Forms
            Start-Sleep -Seconds 2
            [System.Windows.Forms.SendKeys]::SendWait("{TAB}{TAB}{TAB}")
            Start-Sleep -Seconds 1
            [System.Windows.Forms.SendKeys]::SendWait("{ENTER}")
            """
            
            subprocess.run(["powershell", "-Command", ps_command], 
                          timeout=10, capture_output=True)
            
            print("Attempted to navigate to and play first video")
            return True
            
        except Exception as e:
            print(f"PowerShell automation failed: {e}")
        
        # Method 3: Fallback - open search with better parameters using the determined search term
        query = urllib.parse.quote_plus(final_search_term)
        # Use parameters that might make the first video more clickable
        search_url = f"https://www.youtube.com/results?search_query={query}&sp=EgIYAw%253D%253D&tbm=vid"
        webbrowser.open(search_url)
        
        print("Opened YouTube search - the first video should be prominently displayed")
        print("Note: For automatic playing, consider installing selenium and ChromeDriver")
        
        return True
        
    except Exception as e:
        print(f"Error playing YouTube video directly: {e}")
        return False


def play_youtube_video_ultra_direct(search_term: str) -> bool:
    """Ultra-direct YouTube video playback - gets video URL and opens directly without browser automation
    
    This is the most direct method that:
    1. Uses yt-dlp to get direct video URL
    2. Uses YouTube API to extract video ID 
    3. Uses direct video URLs
    4. Bypasses search page entirely
    
    Args:
        search_term: What to search for on YouTube
        
    Returns:
        True if successful, False otherwise
    """
    try:
        import webbrowser
        import urllib.parse
        import subprocess
        import json
        import re
        
        print(f"Attempting ultra-direct playback for: {search_term}")
        
        # Method 1: yt-dlp with improved parameters for direct play
        try:
            # Enhanced yt-dlp command for best results
            result = subprocess.run([
                "yt-dlp", 
                "--get-title",
                "--get-url",
                "--format", "best[height<=720]",  # Get good quality but not too large
                "--no-playlist",  # Only get single video
                "--ignore-errors",
                f"ytsearch1:{search_term}"
            ], capture_output=True, text=True, timeout=15)
            
            if result.returncode == 0 and result.stdout.strip():
                lines = result.stdout.strip().split('\n')
                if len(lines) >= 2:
                    title = lines[0]
                    url = lines[1]
                    print(f"✅ Direct video found: {title}")
                    print(f"🎥 Opening direct URL: {url}")
                    
                    # Open the direct video URL - this should play immediately
                    webbrowser.open(url)
                    return True
                else:
                    print("yt-dlp returned unexpected format")
            else:
                print(f"yt-dlp failed: {result.stderr}")
        except (subprocess.TimeoutExpired, FileNotFoundError, Exception) as e:
            print(f"yt-dlp method failed: {e}")
        
        # Method 2: Enhanced video ID extraction with better regex patterns
        try:
            import requests
            
            query = urllib.parse.quote_plus(search_term)
            search_url = f"https://www.youtube.com/results?search_query={query}"
            
            headers = {
                'User-Agent': 'Mozilla/5.0 (Windows NT 10.0; Win64; x64) AppleWebKit/537.36 (KHTML, like Gecko) Chrome/91.0.4472.124 Safari/537.36'
            }
            
            response = requests.get(search_url, headers=headers, timeout=10)
            if response.status_code == 200:
                # Multiple regex patterns to catch different YouTube HTML formats
                video_id_patterns = [
                    r'"videoId":"([^"]{11})"',  # Standard format
                    r'"videoId":"([^"]+)"',     # Alternative format
                    r'/watch\?v=([^"&\s]+)',    # URL format in HTML
                    r'"url":"/watch\?v=([^"]+)"'  # URL in JSON
                ]
                
                for pattern in video_id_patterns:
                    video_ids = re.findall(pattern, response.text)
                    if video_ids:
                        video_id = video_ids[0]
                        # Ensure it's a valid YouTube video ID (11 characters)
                        if len(video_id) == 11 and video_id.replace('_', '').replace('-', '').isalnum():
                            watch_url = f"https://www.youtube.com/watch?v={video_id}"
                            print(f"✅ Found video ID: {video_id}")
                            print(f"🎥 Opening: {watch_url}")
                            webbrowser.open(watch_url)
                            return True
                
                print("No valid video ID found in search results")
            else:
                print(f"Failed to fetch search results: HTTP {response.status_code}")
                
        except (ImportError, Exception) as e:
            print(f"Video ID extraction failed: {e}")
        
        # Method 3: Use YouTube's direct search API endpoint (if available)
        try:
            import requests
            
            # Alternative approach using YouTube's internal API patterns
            api_url = f"https://www.youtube.com/results?search_query={urllib.parse.quote_plus(search_term)}&sp=EgIYAQ%253D%253D"
            
            headers = {
                'User-Agent': 'Mozilla/5.0 (Windows NT 10.0; Win64; x64) AppleWebKit/537.36',
                'Accept': 'text/html,application/xhtml+xml,application/xml;q=0.9,*/*;q=0.8'
            }
            
            response = requests.get(api_url, headers=headers, timeout=10)
            if response.status_code == 200:
                # Look for ytInitialData containing video information
                initial_data_match = re.search(r'var ytInitialData = ({.+?});', response.text)
                if initial_data_match:
                    try:
                        data = json.loads(initial_data_match.group(1))
                        # Navigate through YouTube's complex JSON structure to find first video
                        contents = data.get('contents', {}).get('twoColumnSearchResultsRenderer', {}).get('primaryContents', {}).get('sectionListRenderer', {}).get('contents', [])
                        
                        for section in contents:
                            items = section.get('itemSectionRenderer', {}).get('contents', [])
                            for item in items:
                                video_renderer = item.get('videoRenderer', {})
                                if video_renderer:
                                    video_id = video_renderer.get('videoId')
                                    if video_id:
                                        watch_url = f"https://www.youtube.com/watch?v={video_id}"
                                        title = video_renderer.get('title', {}).get('runs', [{}])[0].get('text', 'Unknown')
                                        print(f"✅ Found via API: {title}")
                                        print(f"🎥 Opening: {watch_url}")
                                        webbrowser.open(watch_url)
                                        return True
                    except (json.JSONDecodeError, KeyError, IndexError) as e:
                        print(f"Failed to parse YouTube data: {e}")
        
        except (ImportError, Exception) as e:
            print(f"YouTube API method failed: {e}")
        
        # Method 4: Fallback - direct search with optimized URL
        try:
            query = urllib.parse.quote_plus(search_term)
            # Use URL parameters that help with direct video access
            direct_url = f"https://www.youtube.com/results?search_query={query}&sp=EgIYBA%253D%253D&tbm=vid"
            print(f"🔄 Fallback: Opening optimized search URL")
            webbrowser.open(direct_url)
            return True
            
        except Exception as e:
            print(f"Fallback method failed: {e}")
            return False
            
    except Exception as e:
        print(f"Ultra-direct play failed: {e}")
        return False


def auto_click_first_youtube_video() -> bool:
    """Automatically clicks the first YouTube video that appears on screen
    
    This function tries to find and click the first video thumbnail after
    a YouTube search page has loaded.
    
    Returns:
        True if successful, False otherwise
    """
    try:
        import time
        
        # Method 1: Try with pyautogui
        try:
            import pyautogui
            
            # Wait a moment for page to load
            time.sleep(3)
            
            # Take a screenshot to find video thumbnails
            # Look for the distinctive YouTube thumbnail pattern
            screenshot = pyautogui.screenshot()
            
            # Try to find and click the first video thumbnail
            # This is approximate positioning - you may need to adjust
            screen_width, screen_height = pyautogui.size()
            
            # Common position where first video appears after search
            click_positions = [
                (screen_width // 3, screen_height // 3),    # Top-left area
                (screen_width // 2, screen_height // 2),    # Center area  
                (screen_width // 4, screen_height // 4),    # Small offset from top-left
            ]
            
            for x, y in click_positions:
                try:
                    pyautogui.moveTo(x, y, duration=0.5)
                    pyautogui.click()
                    print(f"Clicked at position ({x}, {y})")
                    time.sleep(2)  # Wait to see if video started
                    return True
                except Exception:
                    continue
                    
        except ImportError:
            print("pyautogui not available for auto-clicking")
        except Exception as e:
            print(f"Auto-click with pyautogui failed: {e}")
        
        # Method 2: Try with Windows automation
        try:
            import win32gui
            import win32con
            import win32api
            
            # Find the browser window
            def enum_windows_callback(hwnd, windows):
                if win32gui.IsWindowVisible(hwnd):
                    window_title = win32gui.GetWindowText(hwnd)
                    if "youtube.com" in window_title.lower() or "chrome" in window_title.lower():
                        windows.append(hwnd)
            
            windows = []
            win32gui.EnumWindows(enum_windows_callback, windows)
            
            if windows:
                browser_hwnd = windows[0]
                win32gui.SetForegroundWindow(browser_hwnd)
                
                # Send Tab keys to navigate to first video, then Enter
                win32api.keybd_event(win32con.VK_TAB, 0, 0, 0)
                time.sleep(0.1)
                win32api.keybd_event(win32con.VK_TAB, 0, win32con.KEYEVENTF_KEYUP, 0)
                
                win32api.keybd_event(win32con.VK_TAB, 0, 0, 0)
                time.sleep(0.1)
                win32api.keybd_event(win32con.VK_TAB, 0, win32con.KEYEVENTF_KEYUP, 0)
                
                time.sleep(0.5)
                win32api.keybd_event(win32con.VK_RETURN, 0, 0, 0)
                win32api.keybd_event(win32con.VK_RETURN, 0, win32con.KEYEVENTF_KEYUP, 0)
                
                print("Attempted keyboard navigation to first video")
                return True
                
        except ImportError:
            print("win32gui not available for Windows automation")
        except Exception as e:
            print(f"Windows automation failed: {e}")
        
        return False
        
    except Exception as e:
        print(f"Error in auto_click_first_youtube_video: {e}")
        return False


def skip_youtube_ad() -> bool:
    """Automatically skips the first YouTube ad that appears
    
    This function tries to find and click the "Skip Ad" button or handles
    the first ad that appears when playing a YouTube video.
    
    Returns:
        True if ad was skipped successfully, False otherwise
    """
    try:
        import time
        
        print("Looking for YouTube ad to skip...")
        
        # Method 1: Try with pyautogui to find and click "Skip Ad" button
        try:
            import pyautogui
            
            # Wait a moment for ad to load
            time.sleep(2)
            
            # Get screen dimensions
            screen_width, screen_height = pyautogui.size()
            
            # Common positions where "Skip Ad" button appears (adjust as needed)
            # The button is usually in the bottom-right area of the video player
            skip_ad_positions = [
                (int(screen_width * 0.85), int(screen_height * 0.75)),  # Bottom-right
                (int(screen_width * 0.80), int(screen_height * 0.70)),  # Slightly up-left
                (int(screen_width * 0.90), int(screen_height * 0.80)),  # Far bottom-right
                (int(screen_width * 0.75), int(screen_height * 0.85)),  # Lower-middle-right
            ]
            
            # Also try clicking on center area where "Skip Ad" text might appear
            center_positions = [
                (screen_width // 2, int(screen_height * 0.6)),   # Center area
                (screen_width // 2, int(screen_height * 0.7)),   # Lower center
                (screen_width // 2, int(screen_height * 0.5)),   # Upper center
            ]
            
            all_positions = skip_ad_positions + center_positions
            
            for x, y in all_positions:
                try:
                    pyautogui.moveTo(x, y, duration=0.3)
                    pyautogui.click()
                    time.sleep(1)  # Wait to see if ad was skipped
                    print(f"Attempted to skip ad at position ({x}, {y})")
                    # Return True on any click attempt as we can't reliably verify
                    return True
                except Exception:
                    continue
                    
        except ImportError:
            print("pyautogui not available for ad skipping")
        except Exception as e:
            print(f"Auto-click ad skip failed: {e}")
        
        # Method 2: Try keyboard shortcuts
        try:
            import subprocess
            
            # Use keyboard navigation to try to skip ad
            ps_command = """
            Add-Type -AssemblyName System.Windows.Forms
            Start-Sleep -Seconds 1
            # Try pressing 'k' key (play/pause toggle that sometimes skips ads)
            [System.Windows.Forms.SendKeys]::SendWait("k")
            Start-Sleep -Seconds 0.5
            # Try pressing space (pause/play)
            [System.Windows.Forms.SendKeys]::SendWait(" ")
            Start-Sleep -Seconds 0.5
            # Try pressing right arrow (skip forward)
            [System.Windows.Forms.SendKeys]::SendWait("{RIGHT}")
            """
            
            subprocess.run(["powershell", "-Command", ps_command], 
                          timeout=10, capture_output=True)
            
            print("Attempted keyboard navigation to skip ad")
            return True
            
        except Exception as e:
            print(f"Keyboard ad skip failed: {e}")
        
        # Method 3: Try with Windows API for more precise control
        try:
            import win32gui
            import win32con
            import win32api
            
            # Find the browser window with YouTube
            def enum_windows_callback(hwnd, windows):
                if win32gui.IsWindowVisible(hwnd):
                    window_title = win32gui.GetWindowText(hwnd)
                    if "youtube.com" in window_title.lower() or "chrome" in window_title.lower():
                        windows.append(hwnd)
            
            windows = []
            win32gui.EnumWindows(enum_windows_callback, windows)
            
            if windows:
                browser_hwnd = windows[0]
                win32gui.SetForegroundWindow(browser_hwnd)
                
                # Send 'k' key to toggle play/pause (sometimes skips ads)
                win32api.keybd_event(ord('K'), 0, 0, 0)
                time.sleep(0.1)
                win32api.keybd_event(ord('K'), 0, win32con.KEYEVENTF_KEYUP, 0)
                
                time.sleep(0.5)
                
                # Send right arrow to skip forward
                win32api.keybd_event(win32con.VK_RIGHT, 0, 0, 0)
                win32api.keybd_event(win32con.VK_RIGHT, 0, win32con.KEYEVENTF_KEYUP, 0)
                
                print("Attempted Windows API key press to skip ad")
                return True
                
        except ImportError:
            print("win32gui not available for Windows API ad skipping")
        except Exception as e:
            print(f"Windows API ad skip failed: {e}")
        
        print("Could not automatically skip ad - manual intervention may be required")
        return False
        
    except Exception as e:
        print(f"Error in skip_youtube_ad: {e}")
        return False


def open_youtube_skip_ad_and_play(search_term: str) -> bool:
    """Opens YouTube, plays first video, and automatically skips the first ad
    
    Args:
        search_term: What to search for on YouTube
        
    Returns:
        True if successful, False otherwise
    """
    try:
        print(f"Opening YouTube, playing video, and preparing to skip ads for: {search_term}")
        
        # First open and play the video
        success = open_youtube_and_play_video(search_term)
        
        if success:
            # Wait a moment for the video to start loading
            import time
            time.sleep(3)
            
            # Attempt to skip the first ad
            ad_skipped = skip_youtube_ad()
            
            if ad_skipped:
                print("Successfully handled the first ad")
            else:
                print("Video is playing, but manual ad skipping may be needed")
        
        return success
        
    except Exception as e:
        print(f"Error in open_youtube_skip_ad_and_play: {e}")
        return False


def create_powerpoint_presentation(topic: str, filename: str = None, save_path: str = None) -> bool:
    """
    Creates a PowerPoint presentation on a given topic using python-pptx library
    
    Args:
        topic: The topic for the presentation (e.g., "pollution", "climate change")
        filename: Optional filename for the presentation
        save_path: Optional full path where to save the presentation (default: Desktop)
        
    Returns:
        True if presentation was created successfully, False otherwise
    """
    try:
        import os
        import sys
        
        # Set default filename if not provided
        if not filename:
            # Sanitize topic for filename
            safe_topic = "".join(c for c in topic if c.isalnum() or c in (' ', '-', '_')).strip()
            safe_topic = safe_topic.replace(' ', '_')[:50]  # Limit length
            filename = f"{safe_topic}_Presentation.pptx"
        
        # Determine save path
        if save_path:
            # If user specified save_path, use it (like "D:\\")
            if save_path.endswith('\\'):
                filepath = os.path.join(save_path, filename)
            else:
                filepath = save_path
        else:
            # Default to desktop
            desktop = os.path.join(os.path.expanduser("~"), "Desktop")
            filepath = os.path.join(desktop, filename)
        
        print(f"Creating PowerPoint presentation about '{topic}'...")
        print(f"Save location: {filepath}")
        
        try:
            from pptx import Presentation
            from pptx.enum.text import PP_ALIGN
            from pptx.util import Pt, Inches
            
            # Create presentation object
            prs = Presentation()
            
            # Define content based on topic
            if "pollution" in topic.lower():
                slides_content = [
                    {
                        "title": "Pollution: A Global Crisis",
                        "subtitle": "Understanding the Causes, Effects, and Solutions"
                    },
                    {
                        "title": "Types of Pollution",
                        "content": "• Air Pollution\n• Water Pollution\n• Land Pollution\n• Noise Pollution\n• Light Pollution"
                    },
                    {
                        "title": "Causes of Pollution",
                        "content": "• Industrial Emissions\n• Agricultural Runoff\n• Deforestation\n• Fossil Fuel Combustion\n• Waste Disposal"
                    },
                    {
                        "title": "Effects of Pollution",
                        "content": "• Health Problems\n• Environmental Degradation\n• Climate Change\n• Loss of Biodiversity\n• Economic Impacts"
                    },
                    {
                        "title": "Solutions to Pollution",
                        "content": "• Reduce, Reuse, Recycle\n• Renewable Energy\n• Sustainable Agriculture\n• Pollution Control Technologies\n• Government Regulations"
                    }
                ]
            else:
                # Generic content structure for any topic
                slides_content = [
                    {
                        "title": f"{topic.title()}: Overview",
                        "subtitle": "A comprehensive analysis"
                    },
                    {
                        "title": f"Understanding {topic.title()}",
                        "content": f"• Key concepts and definitions\n• Important factors\n• Current status\n• Future outlook"
                    },
                    {
                        "title": f"Impact of {topic.title()}",
                        "content": "• Environmental effects\n• Social implications\n• Economic considerations\n• Global perspective"
                    },
                    {
                        "title": "Conclusion",
                        "content": f"• Summary of {topic}\n• Key takeaways\n• Next steps\n• Call to action"
                        }
                    ]
            
            # Add slides
            for i, slide_data in enumerate(slides_content):
                try:
                    if i == 0:  # Title slide
                        slide_layout = prs.slide_layouts[0]  # Layout 0 = Title slide
                        slide = prs.slides.add_slide(slide_layout)
                        title = slide.shapes.title
                        subtitle = slide.placeholders[1]
                        
                        title.text = slide_data["title"]
                        subtitle.text = slide_data.get("subtitle", "")
                        
                    else:  # Content slides
                        slide_layout = prs.slide_layouts[1]  # Layout 1 = Title and content
                        slide = prs.slides.add_slide(slide_layout)
                        title = slide.shapes.title
                        content = slide.placeholders[1]
                        
                        title.text = slide_data["title"]
                        tf = content.text_frame
                        tf.text = slide_data["content"]
                        
                    print(f"Added slide {i+1}: {slide_data['title']}")
                    
                except Exception as slide_error:
                    print(f"Warning: Error adding slide {i+1}: {slide_error}")
                    continue
            
            # Save the presentation
            print(f"Saving presentation to: {filepath}")
            prs.save(filepath)
            
            print(f"✅ PowerPoint presentation '{filename}' created successfully!")
            print(f"📁 Location: {filepath}")
            return True
            
        except ImportError:
            print("❌ python-pptx library not found. Installing...")
            try:
                import subprocess
                subprocess.check_call([sys.executable, "-m", "pip", "install", "python-pptx"])
                print("✅ python-pptx installed. Please run the command again.")
                return False
            except Exception as install_error:
                print(f"❌ Failed to install python-pptx: {install_error}")
                return False
                
    except Exception as e:
        print(f"❌ Error creating PowerPoint presentation: {e}")
        return False

def write_text_to_file(content: str, filename: str, save_path: str = None) -> bool:
    """
    Writes text content to a file at the specified location
    
    Args:
        content: The text content to write
        filename: Name of the file to create
        save_path: Optional path where to save the file (default: D drive)
        
    Returns:
        True if file was created successfully, False otherwise
    """
    try:
        import os
        
        # Determine save path
        if save_path:
            # If user specified save_path, use it
            if save_path.endswith('\\'):
                filepath = os.path.join(save_path, filename)
            else:
                filepath = save_path
        else:
            # Default to D drive
            filepath = os.path.join("D:\\", filename)
        
        # Ensure the directory exists
        directory = os.path.dirname(filepath)
        if directory and not os.path.exists(directory):
            try:
                os.makedirs(directory)
                print(f"Created directory: {directory}")
            except Exception as dir_error:
                print(f"Warning: Could not create directory {directory}: {dir_error}")
        
        # Write the file
        print(f"Writing content to: {filepath}")
        with open(filepath, 'w', encoding='utf-8') as f:
            f.write(content)
        
        print(f"✅ File '{filename}' created successfully!")
        print(f"📁 Location: {filepath}")
        return True
        
    except Exception as e:
        print(f"❌ Error writing file: {e}")
        return False

def create_ai_news_file(filename: str = "ai_news.txt") -> bool:
    """Creates a text file with latest AI news, with robust fallback content
    
    Args:
        filename: Name of the file to create (default: "ai_news.txt")
        
    Returns:
        True if file was created successfully, False otherwise
    """
    try:
        import os
        import datetime
        
        # Get desktop path
        desktop = os.path.join(os.path.expanduser("~"), "Desktop")
        filepath = os.path.join(desktop, filename)
        
        # Initialize news content
        timestamp = datetime.datetime.now().strftime('%Y-%m-%d %H:%M:%S')
        news_content = f"Latest AI News - {timestamp}\n\n"
        
        success_count = 0
        total_sources = 0
        
        # Try to scrape news from various sources
        try:
            import requests
            from bs4 import BeautifulSoup
            
            # Define news sources with multiple selector patterns for reliability
            news_sources = [
                {
                    "name": "TechCrunch AI",
                    "url": "https://techcrunch.com/category/artificial-intelligence/",
                    "selectors": [
                        ".post-block__title__link",
                        "h2.post-block__title a",
                        ".river-block__title a",
                        "h2 a",
                        ".headline a",
                        "article h2 a",
                        ".post-title a",
                        "[data-testid='post-title'] a",
                        ".cw9nws a",
                        "h3 a"
                    ]
                },
                {
                    "name": "MIT Technology Review",
                    "url": "https://www.technologyreview.com/topic/artificial-intelligence/",
                    "selectors": [
                        "h3 a",
                        ".articleTitle a",
                        ".post-title a",
                        ".headline a",
                        "article h2 a",
                        ".river-block__headline a",
                        ".post-block__title a",
                        "h2 a",
                        "a[data-testid*='headline']",
                        ".summary a"
                    ]
                },
                {
                    "name": "ArXiv AI Papers",
                    "url": "https://arxiv.org/list/cs.AI/recent",
                    "selectors": [
                        ".list-title math .ltx_Math",
                        ".list-title a",
                        ".list-identifier a",
                        ".arxiv-result a",
                        ".title a",
                        "a[title*='arXiv']",
                        ".list-identifier",
                        ".list-title"
                    ]
                },
                {
                    "name": "AI News (Alternative)",
                    "url": "https://www.artificialintelligence-news.com/",
                    "selectors": [
                        "h2 a",
                        ".entry-title a",
                        ".post-title a",
                        "article h2 a",
                        ".headline a"
                    ]
                },
                {
                    "name": "VentureBeat AI",
                    "url": "https://venturebeat.com/ai/",
                    "selectors": [
                        ".article-primary a",
                        ".headline a",
                        "h2 a",
                        "article h2 a",
                        ".post-title a",
                        ".article-header a"
                    ]
                }
            ]
            
            headers = {
                'User-Agent': 'Mozilla/5.0 (Windows NT 10.0; Win64; x64) AppleWebKit/537.36 (KHTML, like Gecko) Chrome/91.0.4472.124 Safari/537.36',
                'Accept': 'text/html,application/xhtml+xml,application/xml;q=0.9,image/webp,*/*;q=0.8',
                'Accept-Language': 'en-US,en;q=0.5',
                'Accept-Encoding': 'gzip, deflate',
                'DNT': '1',
                'Connection': 'keep-alive',
                'Upgrade-Insecure-Requests': '1'
            }
            
            for source in news_sources:
                total_searches += 1
                try:
                    print(f"Fetching news from {source['name']}...")
                    response = requests.get(source["url"], headers=headers, timeout=15, allow_redirects=True)
                    response.raise_for_status()
                    print(f"Successfully connected to {source['name']} (Status: {response.status_code})")
                    
                    soup = BeautifulSoup(response.content, "html.parser")
                    source_news = []
                    selector_worked = None
                    
                    # Try multiple selectors until we find content
                    for selector in source["selectors"]:
                        try:
                            headlines = soup.select(selector)
                            print(f"  Trying selector '{selector}': found {len(headlines)} elements")
                            
                            if headlines:
                                for headline in headlines[:5]:  # Limit to 5 items per source
                                    try:
                                        title = headline.text.strip()
                                        if title and len(title) > 10:  # Ensure meaningful content
                                            link = headline.get("href", "")
                                            if link and not link.startswith("http"):
                                                if source["url"].endswith("/"):
                                                    link = source["url"] + link
                                                else:
                                                    link = source["url"] + "/" + link
                                            
                                            source_news.append(f"• {title}")
                                            if link and link.startswith("http"):
                                                source_news.append(f"  Link: {link}\n")
                                            else:
                                                source_news.append("")
                                            
                                            if len(source_news) >= 8:  # Limit total items
                                                break
                                    except Exception as item_error:
                                        print(f"    Error processing item: {item_error}")
                                        continue
                                        
                            if source_news:
                                selector_worked = selector
                                print(f"  Successfully extracted content with selector: {selector}")
                                break
                                
                        except Exception as selector_error:
                            print(f"  Error with selector '{selector}': {selector_error}")
                            continue
                    
                    if source_news:
                        news_content += f"=== {source['name']} ===\n"
                        news_content += "\n".join(source_news) + "\n\n"
                        success_count += 1
                        print(f"Successfully fetched {len([x for x in source_news if x.startswith('•')])} items from {source['name']}")
                    else:
                        print(f"No content found for {source['name']} - tried {len(source['selectors'])} selectors")
                        
                except requests.exceptions.RequestException as e:
                    print(f"Network error fetching from {source['name']}: {e}")
                    continue
                except Exception as e:
                    print(f"Unexpected error fetching from {source['name']}: {e}")
                    continue
            
        except ImportError:
            print("requests or BeautifulSoup not available, using fallback content")
        
        # If no news was successfully scraped, add fallback content
        if success_count == 0:
            news_content += """=== AI News Fallback Content ===

Recent AI Developments (as of """ + timestamp + """):

• AI Language Models Continue to Advance
  Major developments in natural language processing and understanding

• Machine Learning in Healthcare
  New applications of AI in medical diagnosis and treatment

• Autonomous Systems Research
  Progress in self-driving vehicles and robotic systems

• AI Ethics and Governance
  Ongoing discussions about responsible AI development

• Computer Vision Breakthroughs
  Advances in image recognition and visual AI systems

Note: This content was generated as fallback since real-time news scraping was unavailable.
For current news, please visit:
- https://techcrunch.com/category/artificial-intelligence/
- https://www.technologyreview.com/topic/artificial-intelligence/
- https://arxiv.org/list/cs.AI/recent

"""
        else:
            news_content += f"\nTotal sources successfully scraped: {success_count}/{total_sources}\n"
            news_content += "Generated: " + timestamp
        
        # Ensure content is not empty before writing
        if len(news_content.strip()) < 50:
            print("Warning: Generated content is very short, adding additional fallback...")
            news_content += "\n\nNote: Limited content was retrieved. Check your internet connection or try again later."
        
        # Write to file
        with open(filepath, "w", encoding="utf-8") as f:
            f.write(news_content)
        
        print(f"AI news file created successfully at: {filepath}")
        print(f"File size: {len(news_content)} characters")
        return True
        
    except Exception as e:
        print(f"Error creating AI news file: {e}")
        return False


def create_news_file(topic: str = "general", filename: str = None) -> bool:
    """Creates a text file with latest news for any topic, with robust fallback content
    
    Args:
        topic: News topic (e.g., "sports", "technology", "politics", "health", etc.)
        filename: Name of the file to create (default: "{topic}_news.txt")
        
    Returns:
        True if file was created successfully, False otherwise
    """
    try:
        import os
        import datetime
        
        # Set default filename if not provided
        if not filename:
            filename = f"{topic.replace(' ', '_')}_news.txt"
        
        # Get desktop path
        desktop = os.path.join(os.path.expanduser("~"), "Desktop")
        filepath = os.path.join(desktop, filename)
        
        # Initialize news content
        timestamp = datetime.datetime.now().strftime('%Y-%m-%d %H:%M:%S')
        news_content = f"Latest {topic.title()} News - {timestamp}\n\n"
        
        success_count = 0
        total_sources = 0
        
        # Define news sources based on topic
        news_sources = _get_news_sources_for_topic(topic)
        
        # Try to scrape news from various sources
        try:
            import requests
            from bs4 import BeautifulSoup
            
            headers = {
                'User-Agent': 'Mozilla/5.0 (Windows NT 10.0; Win64; x64) AppleWebKit/537.36 (KHTML, like Gecko) Chrome/91.0.4472.124 Safari/537.36'
            }
            
            for source in news_sources:
                total_sources += 1
                try:
                    print(f"Fetching {topic} news from {source['name']}...")
                    response = requests.get(source["url"], headers=headers, timeout=10)
                    response.raise_for_status()
                    
                    soup = BeautifulSoup(response.content, "html.parser")
                    source_news = []
                    
                    # Try multiple selectors until we find content
                    for selector in source["selectors"]:
                        headlines = soup.select(selector)
                        if headlines:
                            for headline in headlines[:5]:  # Limit to 5 items per source
                                try:
                                    title = headline.text.strip()
                                    if title and len(title) > 10:  # Ensure meaningful content
                                        link = headline.get("href", "")
                                        if link and not link.startswith("http"):
                                            if source["url"].endswith("/"):
                                                link = source["url"] + link
                                            else:
                                                link = source["url"] + "/" + link
                                        
                                        source_news.append(f"• {title}")
                                        if link and link.startswith("http"):
                                            source_news.append(f"  Link: {link}\n")
                                        else:
                                            source_news.append("")
                                        
                                        if len(source_news) >= 8:  # Limit total items
                                            break
                                except:
                                    continue
                            if source_news:
                                break
                    
                    if source_news:
                        news_content += f"=== {source['name']} ===\n"
                        news_content += "\n".join(source_news) + "\n\n"
                        success_count += 1
                        print(f"Successfully fetched {len([x for x in source_news if x.startswith('•')])} items from {source['name']}")
                    else:
                        print(f"No content found for {source['name']}")
                        
                except Exception as e:
                    print(f"Error fetching from {source['name']}: {e}")
                    continue
            
        except ImportError:
            print("requests or BeautifulSoup not available, using fallback content")
        
        # If no news was successfully scraped, add fallback content
        if success_count == 0:
            news_content += _get_fallback_content_for_topic(topic, timestamp)
        else:
            news_content += f"\nTotal sources successfully scraped: {success_count}/{total_sources}\n"
            news_content += "Generated: " + timestamp
        
        # Ensure content is not empty before writing
        if len(news_content.strip()) < 50:
            print("Warning: Generated content is very short, adding additional fallback...")
            news_content += f"\n\nNote: Limited content was retrieved for {topic} news. Check your internet connection or try again later."
        
        # Write to file
        with open(filepath, "w", encoding="utf-8") as f:
            f.write(news_content)
        
        print(f"{topic.title()} news file created successfully at: {filepath}")
        print(f"File size: {len(news_content)} characters")
        return True
        
    except Exception as e:
        print(f"Error creating {topic} news file: {e}")
        return False


def _get_news_sources_for_topic(topic: str) -> list:
    """Get appropriate news sources for a given topic"""
    topic_lower = topic.lower()
    
    # Define base news sources that work for most topics
    base_sources = [
        {
            "name": "BBC News",
            "url": f"https://www.bbc.com/news/{topic_lower.replace(' ', '-')}",
            "selectors": [
                ".gs-c-promo-heading__title",
                "h3.gs-c-promo-heading__title",
                ".gs-c-promo a",
                "h2 a"
            ]
        },
        {
            "name": "CNN",
            "url": f"https://www.cnn.com/{topic_lower.replace(' ', '/')}",
            "selectors": [
                ".cd__headline-text",
                "h3.cd__headline-text",
                ".zn-body__read-all a",
                ".cd a"
            ]
        },
        {
            "name": "Reuters",
            "url": f"https://www.reuters.com/{topic_lower.replace(' ', '/')}",
            "selectors": [
                "[data-testid='Heading']",
                ".media-story-card__headline",
                "h3 a",
                ".story-collection__wrapper a"
            ]
        }
    ]
    
    # Add topic-specific sources
    if "tech" in topic_lower or "technology" in topic_lower:
        base_sources.extend([
            {
                "name": "TechCrunch",
                "url": f"https://techcrunch.com/category/{topic_lower.replace(' ', '-')}/",
                "selectors": [".post-block__title__link", "h2.post-block__title a", ".river-block__title a"]
            }
        ])
    elif "sport" in topic_lower:
        base_sources.extend([
            {
                "name": "ESPN",
                "url": f"https://www.espn.com/{topic_lower.replace(' ', '/')}/",
                "selectors": [".contentItem__title", "h1 a", ".headlineStack__list a"]
            }
        ])
    elif "health" in topic_lower or "medical" in topic_lower:
        base_sources.extend([
            {
                "name": "WebMD",
                "url": f"https://www.webmd.com/news/{topic_lower.replace(' ', '-')}/",
                "selectors": [".headline", "h2 a", ".headlines-list a"]
            }
        ])
    
    return base_sources


def _get_fallback_content_for_topic(topic: str, timestamp: str) -> str:
    """Generate fallback content for a specific news topic"""
    topic_capitalized = topic.title()
    
    fallback_templates = {
        "sports": f"""=== {topic_capitalized} News Fallback Content ===

Recent Sports Developments (as of {timestamp}):

• Major League Updates
  Latest scores, standings, and player statistics

• Transfer News and Rumors
  Updates on player movements and contract negotiations

• Championship Races
  Current standings in major tournaments and leagues

• Injury Reports
  Updates on player health and recovery status

• Olympic and International Events
  Preparations and results from global competitions

Note: This content was generated as fallback since real-time news scraping was unavailable.
For current sports news, please visit:
- https://www.espn.com/
- https://www.bbc.com/sport
- https://www.reuters.com/sports/

""",
        "technology": f"""=== {topic_capitalized} News Fallback Content ===

Recent Technology Developments (as of {timestamp}):

• Software and App Updates
  Latest releases and feature announcements

• Hardware Innovations
  New devices, processors, and tech infrastructure

• Cybersecurity Alerts
  Important security updates and threat intelligence

• Industry Mergers and Acquisitions
  Major business movements in tech sector

• Research and Development
  Breakthrough technologies and scientific advances

Note: This content was generated as fallback since real-time news scraping was unavailable.
For current technology news, please visit:
- https://techcrunch.com/
- https://www.technologyreview.com/
- https://www.wired.com/

""",
        "default": f"""=== {topic_capitalized} News Fallback Content ===

Recent {topic_capitalized} Developments (as of {timestamp}):

• Breaking News Updates
  Latest developments and significant events

• Policy and Regulatory Changes
  Updates on laws, regulations, and governance

• Market and Economic Trends
  Analysis of current market conditions and forecasts

• International Relations
  Global developments and diplomatic updates

• Research and Studies
  Latest findings and research publications

Note: This content was generated as fallback since real-time news scraping was unavailable.
For current {topic_lower} news, please visit:
- https://www.bbc.com/news
- https://www.cnn.com/
- https://www.reuters.com/

"""
    }
    
    # Find appropriate template
    topic_lower = topic.lower()
    for key, template in fallback_templates.items():
        if key != "default" and key in topic_lower:
            return template
    
    return fallback_templates["default"]


def scrape_info_about(search_term: str, info_type: str = "general", filename: str = None) -> bool:
    """Comprehensive web scraping function that can find information about any person or company
    
    Args:
        search_term: The name of person/company or topic to search for
        info_type: Type of information to gather ("person", "company", "general", "news", "wiki")
        filename: Optional custom filename (default: auto-generated)
        
    Returns:
        True if file was created successfully, False otherwise
    """
    try:
        import os
        import datetime
        import urllib.parse
        
        # Generate filename if not provided
        if not filename:
            safe_name = "".join(c for c in search_term if c.isalnum() or c in (' ', '-', '_')).rstrip()
            safe_name = safe_name.replace(' ', '_')[:50]  # Limit length
            filename = f"{safe_name}_{info_type}_info.txt"
        
        # Get desktop path
        desktop = os.path.join(os.path.expanduser("~"), "Desktop")
        filepath = os.path.join(desktop, filename)
        
        # Initialize content
        timestamp = datetime.datetime.now().strftime('%Y-%m-%d %H:%M:%S')
        content = f"Information about: {search_term}\n"
        content += f"Search Type: {info_type.title()}\n"
        content += f"Generated: {timestamp}\n\n"
        
        # URL encode the search term
        encoded_term = urllib.parse.quote_plus(search_term)
        
        success_count = 0
        total_searches = 0
        
        # Define sources based on info type
        sources = _get_info_sources(search_term, info_type, encoded_term)
        
        try:
            import requests
            from bs4 import BeautifulSoup
            
            headers = {
                'User-Agent': 'Mozilla/5.0 (Windows NT 10.0; Win64; x64) AppleWebKit/537.36 (KHTML, like Gecko) Chrome/91.0.4472.124 Safari/537.36',
                'Accept': 'text/html,application/xhtml+xml,application/xml;q=0.9,image/webp,*/*;q=0.8',
                'Accept-Language': 'en-US,en;q=0.5',
                'Accept-Encoding': 'gzip, deflate',
                'DNT': '1',
                'Connection': 'keep-alive',
                'Upgrade-Insecure-Requests': '1'
            }
            
            for source in sources:
                total_searches += 1
                source_name = source["name"]
                source_url = source["url"]
                selectors = source["selectors"]
                
                try:
                    print(f"Searching {source_name} for information about '{search_term}'...")
                    response = requests.get(source_url, headers=headers, timeout=15, allow_redirects=True)
                    response.raise_for_status()
                    
                    soup = BeautifulSoup(response.content, "html.parser")
                    source_content = []
                    
                    # Try multiple selectors to extract information
                    for selector_config in selectors:
                        selector = selector_config["selector"]
                        text_type = selector_config.get("type", "content")
                        
                        try:
                            elements = soup.select(selector)
                            if elements:
                                for element in elements[:10]:  # Limit to 10 items per selector
                                    try:
                                        if text_type == "title":
                                            text = element.get('title', '') or element.text.strip()
                                        elif text_type == "link":
                                            text = element.get('href', '')
                                            if text and not text.startswith('http'):
                                                text = urllib.parse.urljoin(source_url, text)
                                        else:
                                            text = element.text.strip()
                                        
                                        if text and len(text) > 5:
                                            # Clean up the text
                                            text = ' '.join(text.split())  # Remove extra whitespace
                                            
                                            if text_type == "link":
                                                source_content.append(f"• Related Link: {text}")
                                            elif text_type == "title":
                                                source_content.append(f"• {text}")
                                            else:
                                                source_content.append(f"• {text}")
                                            
                                            # Add link if available
                                            if text_type != "link":
                                                href = element.get('href')
                                                if href:
                                                    if not href.startswith('http'):
                                                        href = urllib.parse.urljoin(source_url, href)
                                                    source_content.append(f"  Link: {href}")
                                                source_content.append("")
                                                
                                        if len(source_content) >= 20:  # Limit total items per source
                                            break
                                    except Exception as e:
                                        continue
                                    if len(source_content) >= 20:
                                        break
                                        
                                if source_content:
                                    break
                        except Exception as e:
                            continue
                    
                    if source_content:
                        content += f"=== {source_name} ===\n"
                        content += "\n".join(source_content) + "\n\n"
                        success_count += 1
                        print(f"Successfully found {len([x for x in source_content if x.startswith('•')])} items from {source_name}")
                    else:
                        print(f"No specific content found for '{search_term}' on {source_name}")
                        
                except Exception as e:
                    print(f"Error searching {source_name}: {e}")
                    continue
            
        except ImportError:
            print("requests or BeautifulSoup not available, using fallback content")
        
        # If no content was found, add fallback information
        if success_count == 0:
            content += _get_fallback_info(search_term, info_type, timestamp)
        else:
            content += f"\n=== Summary ===\n"
            content += f"Successfully searched {success_count}/{total_searches} sources\n"
            content += f"Generated: {timestamp}\n"
            content += f"For more current information, visit the sources directly or search '{search_term}' online.\n"
        
        # Ensure content is not empty
        if len(content.strip()) < 100:
            content += f"\n\nNote: Limited information was found for '{search_term}'. "
            content += f"This could be due to:\n"
            content += f"- The term not being well-known or having limited online presence\n"
            content += f"- Network connectivity issues\n"
            content += f"- Changes in website structures\n"
            content += f"\nTry searching manually on:\n"
            content += f"- Google: https://www.google.com/search?q={encoded_term}\n"
            content += f"- Wikipedia: https://en.wikipedia.org/wiki/{encoded_term.replace('+', '_')}\n"
        
        # Write to file
        with open(filepath, "w", encoding="utf-8") as f:
            f.write(content)
        
        print(f"Information file created successfully at: {filepath}")
        print(f"File size: {len(content)} characters")
        return True
        
    except Exception as e:
        print(f"Error creating info file for '{search_term}': {e}")
        return False


def _get_info_sources(search_term: str, info_type: str, encoded_term: str) -> list:
    """Get appropriate sources for searching information about a person/company"""
    
    sources = []
    
    # Google Search results - Updated selectors for better compatibility
    sources.append({
        "name": "Google Search Results",
        "url": f"https://www.google.com/search?q={encoded_term}",
        "selectors": [
            {"selector": "h3", "type": "title"},
            {"selector": "h3 span", "type": "title"},
            {"selector": ".yuRUbf a", "type": "link"},
            {"selector": ".yuRUbf a h3", "type": "content"},
            {"selector": ".VwiC3b", "type": "content"},
            {"selector": ".s3v9rd", "type": "content"},
            {"selector": ".g .yuRUbf", "type": "content"},
            {"selector": ".rc .r", "type": "content"},
            {"selector": ".LC20lb", "type": "content"},
            {"selector": ".DKV0Md", "type": "content"}
        ]
    })
    
    # Wikipedia - Use search first, then direct page if available
    sources.append({
        "name": "Wikipedia Search",
        "url": f"https://en.wikipedia.org/wiki/Special:Search?search={encoded_term}",
        "selectors": [
            {"selector": ".mw-search-result-heading a", "type": "content"},
            {"selector": ".mw-search-result-snippet", "type": "content"},
            {"selector": ".mw-search-results li", "type": "content"},
            {"selector": "h1.firstHeading", "type": "content"},
            {"selector": ".mw-parser-output p", "type": "content"},
            {"selector": ".infobox .infobox-data", "type": "content"}
        ]
    })
    
    # Also try direct Wikipedia page as fallback
    sources.append({
        "name": "Wikipedia Direct",
        "url": f"https://en.wikipedia.org/wiki/{encoded_term.replace('+', '_')}",
        "selectors": [
            {"selector": "h1.firstHeading", "type": "content"},
            {"selector": ".mw-parser-output p", "type": "content"},
            {"selector": ".infobox .infobox-data", "type": "content"},
            {"selector": ".sidebar .sidebar-content", "type": "content"}
        ]
    })
    
    # LinkedIn for people and companies
    if info_type in ["person", "company"]:
        sources.append({
            "name": "LinkedIn Search",
            "url": f"https://www.linkedin.com/search/results/all/?keywords={encoded_term}",
            "selectors": [
                {"selector": ".entity-result__title-text a", "type": "content"},
                {"selector": ".entity-result__primary-subtitle", "type": "content"},
                {"selector": ".entity-result__summary", "type": "content"}
            ]
        })
    
    # Company-specific sources
    if info_type == "company":
        # Crunchbase for company info
        sources.append({
            "name": "Crunchbase",
            "url": f"https://www.crunchbase.com/discover/organization/{encoded_term}",
            "selectors": [
                {"selector": ".company-name", "type": "content"},
                {"selector": ".description", "type": "content"},
                {"selector": ".field-label", "type": "content"},
                {"selector": ".field-value", "type": "content"}
            ]
        })
        
        # Yahoo Finance
        sources.append({
            "name": "Yahoo Finance",
            "url": f"https://finance.yahoo.com/quote/{encoded_term}",
            "selectors": [
                {"selector": "h1[data-symbol]", "type": "content"},
                {"selector": ".company-name", "type": "content"},
                {"selector": ".quote-summary", "type": "content"},
                {"selector": ".general-description", "type": "content"}
            ]
        })
    
    # News sources
    if info_type in ["person", "company", "news"]:
        # BBC News
        sources.append({
            "name": "BBC News Search",
            "url": f"https://www.bbc.com/search?q={encoded_term}",
            "selectors": [
                {"selector": ".search-result__headline", "type": "content"},
                {"selector": ".search-result__summary", "type": "content"},
                {"selector": ".search-result__date", "type": "content"}
            ]
        })
        
        # Reuters
        sources.append({
            "name": "Reuters Search",
            "url": f"https://www.reuters.com/search/news?blob={encoded_term}",
            "selectors": [
                {"selector": ".search-result-title", "type": "content"},
                {"selector": ".search-result-excerpt", "type": "content"},
                {"selector": "[data-testid='Heading']", "type": "content"}
            ]
        })
    
    return sources


def _get_fallback_info(search_term: str, info_type: str, timestamp: str) -> str:
    """Generate fallback information when scraping fails"""
    
    fallback_info = f"""=== Information Not Found Online ===

Search conducted for: {search_term}
Type: {info_type.title()}
Date: {timestamp}

The automated search did not return specific information about '{search_term}'. 
This could be because:

1. Limited online presence or publicity
2. Private individual or confidential company information
3. Very recent or very obscure subject
4. Website structure changes affecting automated searches

Recommended manual search sources:
- Google: https://www.google.com/search?q={search_term.replace(' ', '+')}
- Wikipedia: https://en.wikipedia.org/wiki/{search_term.replace(' ', '_')}
- LinkedIn: https://www.linkedin.com/search/results/all/?keywords={search_term.replace(' ', '+')}

For companies specifically, try:
- Company website: https://www.{search_term.lower().replace(' ', '')}.com
- Crunchbase: https://www.crunchbase.com/discover/organization/{search_term.replace(' ', '+')}
- Yahoo Finance: https://finance.yahoo.com/quote/{search_term.replace(' ', '+')}

For people specifically, try:
- Social media platforms (Twitter, LinkedIn, Facebook)
- Professional networks and academic profiles
- News archives and press releases

Note: Always verify information from multiple sources for accuracy.
"""
    
    return fallback_info


def scrape_info_content(search_term: str, info_type: str = "general") -> str:
    """Scrapes information about a person or company and returns the content as a string
    
    Args:
        search_term: The name of person/company or topic to search for
        info_type: Type of information to gather ("person", "company", "general", "news", "wiki")
        
    Returns:
        String containing the scraped information content
    """
    try:
        import datetime
        import urllib.parse
        
        # Initialize content
        timestamp = datetime.datetime.now().strftime('%Y-%m-%d %H:%M:%S')
        content = f"Information about: {search_term}\n"
        content += f"Search Type: {info_type.title()}\n"
        content += f"Generated: {timestamp}\n\n"
        
        # URL encode the search term
        encoded_term = urllib.parse.quote_plus(search_term)
        
        success_count = 0
        total_searches = 0
        
        # Define sources based on info type
        sources = _get_info_sources(search_term, info_type, encoded_term)
        
        try:
            import requests
            from bs4 import BeautifulSoup
            
            headers = {
                'User-Agent': 'Mozilla/5.0 (Windows NT 10.0; Win64; x64) AppleWebKit/537.36 (KHTML, like Gecko) Chrome/91.0.4472.124 Safari/537.36',
                'Accept': 'text/html,application/xhtml+xml,application/xml;q=0.9,image/webp,*/*;q=0.8',
                'Accept-Language': 'en-US,en;q=0.5',
                'Accept-Encoding': 'gzip, deflate',
                'DNT': '1',
                'Connection': 'keep-alive',
                'Upgrade-Insecure-Requests': '1'
            }
            
            for source in sources:
                total_searches += 1
                source_name = source["name"]
                source_url = source["url"]
                selectors = source["selectors"]
                
                try:
                    print(f"Searching {source_name} for information about '{search_term}'...")
                    response = requests.get(source_url, headers=headers, timeout=15, allow_redirects=True)
                    response.raise_for_status()
                    
                    soup = BeautifulSoup(response.content, "html.parser")
                    source_content = []
                    
                    # Try multiple selectors to extract information
                    for selector_config in selectors:
                        selector = selector_config["selector"]
                        text_type = selector_config.get("type", "content")
                        
                        try:
                            elements = soup.select(selector)
                            if elements:
                                for element in elements[:10]:  # Limit to 10 items per selector
                                    try:
                                        if text_type == "title":
                                            text = element.get('title', '') or element.text.strip()
                                        elif text_type == "link":
                                            text = element.get('href', '')
                                            if text and not text.startswith('http'):
                                                text = urllib.parse.urljoin(source_url, text)
                                        else:
                                            text = element.text.strip()
                                        
                                        if text and len(text) > 5:
                                            # Clean up the text
                                            text = ' '.join(text.split())  # Remove extra whitespace
                                            
                                            if text_type == "link":
                                                source_content.append(f"• Related Link: {text}")
                                            elif text_type == "title":
                                                source_content.append(f"• {text}")
                                            else:
                                                source_content.append(f"• {text}")
                                            
                                            # Add link if available
                                            if text_type != "link":
                                                href = element.get('href')
                                                if href:
                                                    if not href.startswith('http'):
                                                        href = urllib.parse.urljoin(source_url, href)
                                                    source_content.append(f"  Link: {href}")
                                                source_content.append("")
                                                
                                        if len(source_content) >= 20:  # Limit total items per source
                                            break
                                    except Exception as e:
                                        continue
                                    if len(source_content) >= 20:
                                        break
                                        
                                if source_content:
                                    break
                        except Exception as e:
                            continue
                    
                    if source_content:
                        content += f"=== {source_name} ===\n"
                        content += "\n".join(source_content) + "\n\n"
                        success_count += 1
                        print(f"Successfully found {len([x for x in source_content if x.startswith('•')])} items from {source_name}")
                    else:
                        print(f"No specific content found for '{search_term}' on {source_name}")
                        
                except Exception as e:
                    print(f"Error searching {source_name}: {e}")
                    continue
            
        except ImportError:
            print("requests or BeautifulSoup not available, using fallback content")
        
        # If no content was found, add fallback information
        if success_count == 0:
            content += _get_fallback_info(search_term, info_type, timestamp)
        else:
            content += f"\n=== Summary ===\n"
            content += f"Successfully searched {success_count}/{total_searches} sources\n"
            content += f"Generated: {timestamp}\n"
            content += f"For more current information, visit the sources directly or search '{search_term}' online.\n"
        
        # Ensure content is not empty
        if len(content.strip()) < 100:
            content += f"\n\nNote: Limited information was found for '{search_term}'. "
            content += f"This could be due to:\n"
            content += f"- The term not being well-known or having limited online presence\n"
            content += f"- Network connectivity issues\n"
            content += f"- Changes in website structures\n"
            content += f"\nTry searching manually on:\n"
            content += f"- Google: https://www.google.com/search?q={encoded_term}\n"
            content += f"- Wikipedia: https://en.wikipedia.org/wiki/{encoded_term.replace('+', '_')}\n"
        
        return content
        
    except Exception as e:
        error_content = f"Error creating info for '{search_term}': {e}\n"
        error_content += "Please try again or search manually."
        return error_content


# Function mapping for natural language commands - defined after all functions
FUNCTION_MAPPINGS = {
    # Audio/Volume commands
    "mute": mute_system_volume,
    "mute volume": mute_system_volume,
    "mute system": mute_system_volume,
    "unmute": unmute_system_volume,
    "unmute volume": unmute_system_volume,
    "unmute system": unmute_system_volume,
    "set volume": set_system_volume,
    "change volume": set_system_volume,
    "volume": set_system_volume,
    "get volume": get_current_volume,
    "current volume": get_current_volume,
    "check volume": get_current_volume,

    # Desktop commands
    "hide desktop icons": hide_desktop_icons,
    "hide icons": hide_desktop_icons,
    "show desktop icons": show_desktop_icons,
    "show icons": show_desktop_icons,
    "toggle desktop icons": lambda: show_desktop_icons() if not get_desktop_icons_visible() else hide_desktop_icons(),

    # Brightness commands
    "increase brightness": lambda: adjust_brightness(10),
    "decrease brightness": lambda: adjust_brightness(-10),
    "set brightness": set_brightness,
    "brightness": set_brightness,

    # Night light commands
    "night light on": lambda: toggle_night_light(True),
    "night light off": lambda: toggle_night_light(False),
    "toggle night light": toggle_night_light,

    # Airplane mode commands
    "airplane mode on": lambda: toggle_airplane_mode_advanced(True),
    "airplane mode off": lambda: toggle_airplane_mode_advanced(False),
    "toggle airplane mode": toggle_airplane_mode_advanced,

    # System commands
    "restart explorer": restart_explorer,
    "check admin": is_admin,
    "system capabilities": get_system_capabilities,
    "create desktop shortcut": create_desktop_shortcut,
    "create startup shortcut": create_startup_shortcut,
    "remove startup shortcut": remove_startup_shortcut,

    # Camera and Photo commands
    "open camera": open_camera_app,
    "camera": open_camera_app,
    "open camera app": open_camera_app,
    "take screenshot": take_screenshot,
    "screenshot": take_screenshot,
    "capture screen": take_screenshot,
    "open photos": open_photos_app,
    "photos app": open_photos_app,
    
    # YouTube commands
    "play youtube video": open_youtube_and_play_video,
    "youtube play": open_youtube_and_play_video,
    "open youtube and play": open_youtube_and_play_video,
    "play youtube direct": play_youtube_video_ultra_direct,
    "youtube direct": play_youtube_video_ultra_direct,
    "play video direct": play_youtube_video_ultra_direct,
    "direct youtube play": play_youtube_video_ultra_direct,
    "play first youtube video": auto_click_first_youtube_video,
    "skip youtube ad": skip_youtube_ad,
    "skip first ad": skip_youtube_ad,
    "youtube skip ad": skip_youtube_ad,
    "play youtube skip ad": open_youtube_skip_ad_and_play,
    
    # News and content creation commands
    "create ai news": create_ai_news_file,
    "create ai news file": create_ai_news_file,
    "latest ai news": create_ai_news_file,
    "ai news": create_ai_news_file,
    
    # Generic news commands (will be handled by create_news_file function)
    "create news": lambda: create_news_file("general"),
    "create news file": lambda: create_news_file("general"),
    "create sports news": lambda: create_news_file("sports"),
    "create tech news": lambda: create_news_file("technology", "technology_news.txt"),
    "create technology news": lambda: create_news_file("technology"),
    "create health news": lambda: create_news_file("health"),
    "create business news": lambda: create_news_file("business"),
    "create world news": lambda: create_news_file("world"),
    
    # PowerPoint presentation commands
    "create ppt": "create_powerpoint_presentation_function",
    "create powerpoint": "create_powerpoint_presentation_function", 
    "create presentation": "create_powerpoint_presentation_function",
    "make ppt": "create_powerpoint_presentation_function",
    "make powerpoint": "create_powerpoint_presentation_function",
    "ppt based on": "create_powerpoint_presentation_function",
    "powerpoint about": "create_powerpoint_presentation_function",
    "presentation about": "create_powerpoint_presentation_function",
    
    # Information scraping commands (will use scrape_info_about with parameter extraction)
    "search for": "scrape_info_about_function",
    "find information about": "scrape_info_about_function",
    "scrape information about": "scrape_info_about_function",
    "get info about": "scrape_info_about_function",
    "look up": "scrape_info_about_function",
    
    # Text file creation commands
    "write": "write_text_function",
    "write to file": "write_text_function",
    "create text file": "write_text_function",
}


def _parse_scraping_command(command: str):
    """Extract search term and info type from scraping commands"""
    command_lower = command.lower()
    
    # Define scraping trigger phrases and their info types
    trigger_phrases = [
        ("search for", "general"),
        ("find information about", "general"),
        ("scrape information about", "general"),
        ("get info about", "general"),
        ("look up", "general"),
        ("search person", "person"),
        ("search company", "company"),
        ("company info", "company"),
        ("person info", "person"),
        ("news about", "news"),
        ("wikipedia", "wiki")
    ]
    
    # Find the trigger phrase and remove it
    search_term = None
    info_type = "general"
    
    for phrase, detected_type in trigger_phrases:
        if phrase in command_lower:
            # Extract the term after the phrase
            parts = command_lower.split(phrase, 1)
            if len(parts) > 1:
                search_term = parts[1].strip()
                info_type = detected_type
                break
    
    # If no trigger phrase found, try to guess from the command
    if not search_term:
        # Look for common patterns
        if " about " in command_lower:
            search_term = command_lower.split(" about ", 1)[1].strip()
        elif command_lower.startswith(("find ", "search ", "lookup ", "look up ")):
            # Take everything after the initial verb
            for start_word in ["find ", "search ", "lookup ", "look up "]:
                if command_lower.startswith(start_word):
                    search_term = command_lower[len(start_word):].strip()
                    break
    
    # Clean up the search term
    if search_term:
        # Remove extra words and clean up
        search_term = search_term.replace(" and create file", "").replace(" and save", "").replace(" and write to file", "").strip()
        # Take only the first reasonable length (probably a name/topic)
        search_term = " ".join(search_term.split()[:10])  # Limit to 10 words
    
    return search_term, info_type


def _parse_write_command(command: str):
    """Extract content, filename and save path from write commands"""
    import re
    
    command_lower = command.lower()
    
    # Extract filename - look for patterns like "robots.txt" or "to [filename]"
    filename = None
    save_path = None
    content = None
    
    # Look for "to [filename]" pattern
    to_match = re.search(r'to (\S+\.(txt|doc|docx|md))', command_lower)
    if to_match:
        filename = to_match.group(1)
    
    # Look for "in d drive" or "to d:" pattern
    if " in d drive" in command_lower or " to d drive" in command_lower:
        save_path = "D:\\"
    elif " in c drive" in command_lower or " to c drive" in command_lower:
        save_path = "C:\\"
    elif " save to " in command_lower:
        save_to_match = re.search(r"save to\s+([a-z]:\\)", command_lower)
        if save_to_match:
            save_path = save_to_match.group(1)
    
    # Extract the content - everything between "write" and "to"
    # For "write a 200 words paragraph on robots"
    if "paragraph" in command_lower or "on " in command_lower:
        # This is a content generation request
        paragraph_match = re.search(r'(\d+)\s*words?\s*paragraph\s*(?:on|about)\s*(.+)', command_lower)
        if paragraph_match:
            word_count = paragraph_match.group(1)
            topic = paragraph_match.group(2).split()[0]  # First word after "on/about"
            # Generate content (simplified - would use AI in real scenario)
            content = _generate_paragraph_about_topic(topic, int(word_count))
    
    # Extract remaining words as topic/content if content not set
    if not content and " on " in command_lower:
        parts = command_lower.split(" on ")
        if len(parts) > 1:
            content = _generate_paragraph_about_topic(parts[1].split(" to")[0], 200)
    
    # Default filename if not specified
    if not filename and content:
        filename = "content.txt"
    
    return content, filename, save_path

def _generate_paragraph_about_topic(topic: str, word_count: int = 200) -> str:
    """Generate a paragraph about a given topic"""
    # This is a simplified version - in reality this would use AI
    topic_clean = topic.strip()
    
    # Pre-defined content for common topics
    content_templates = {
        "robots": """Robots are increasingly becoming integral parts of modern society, transforming industries and daily life in profound ways. These automated machines, capable of carrying out complex tasks with precision and efficiency, have revolutionized manufacturing, healthcare, exploration, and countless other fields. In factories around the world, robotic arms assemble products with remarkable speed and accuracy, while surgical robots assist doctors in performing delicate procedures with enhanced precision. Exploration robots venture into environments too dangerous for humans, from the depths of the ocean to the surface of Mars, collecting invaluable data and expanding our understanding of the universe. As artificial intelligence continues to evolve, robots are becoming increasingly autonomous and capable of making decisions in real-time. The future promises even more sophisticated robots that can collaborate seamlessly with humans, performing tasks ranging from household chores to complex scientific research. However, this rapid advancement also raises important questions about job displacement, ethics, and the role of robots in society. Regardless of these challenges, robots represent one of humanity's most ambitious technological achievements, pushing the boundaries of what machines can accomplish.""",
    }
    
    if topic_clean.lower() in content_templates:
        return content_templates[topic_clean.lower()]
    
    # Generic template for other topics
    return f"{topic_clean.title()} is a fascinating subject that has captivated human interest for generations. This field encompasses various aspects, from fundamental principles to advanced applications. The study of {topic_clean} involves understanding its core concepts, historical development, and contemporary significance. Researchers and practitioners continue to explore new dimensions of {topic_clean}, uncovering insights that shape our understanding and drive innovation. The impact of {topic_clean} extends across multiple disciplines, influencing technology, society, and the way we perceive the world around us. As we delve deeper into this subject, we discover the interconnectedness of ideas and the complexity of systems involved. The future of {topic_clean} holds promise for groundbreaking discoveries that will further enhance our knowledge and capabilities in this domain."

def _parse_powerpoint_command(command: str):
    """Extract topic and save path from PowerPoint creation commands"""
    command_lower = command.lower()
    
    # Define PowerPoint trigger phrases
    trigger_phrases = [
        ("create ppt based on", ""),
        ("create powerpoint based on", ""),
        ("create presentation based on", ""),
        ("ppt based on", ""),
        ("powerpoint about", ""),
        ("presentation about", ""),
        ("create ppt about", ""),
        ("create powerpoint about", ""),
        ("make ppt about", ""),
        ("make powerpoint about", ""),
        ("generate ppt on", ""),
        ("generate powerpoint on", ""),
    ]
    
    topic = None
    
    for phrase, _ in trigger_phrases:
        if phrase in command_lower:
            # Extract the topic after the phrase
            parts = command_lower.split(phrase, 1)
            if len(parts) > 1:
                topic = parts[1].strip()
                break
    
    # If no trigger phrase found, try other patterns
    if not topic:
        if " about " in command_lower and ("ppt" in command_lower or "powerpoint" in command_lower or "presentation" in command_lower):
            topic = command_lower.split(" about ", 1)[1].strip()
        elif "generate ppt on" in command_lower or "generate powerpoint on" in command_lower:
            parts = command_lower.split(" on ", 1)
            if len(parts) > 1:
                topic = parts[1].strip()
    
    # Extract save location (D drive, C drive, etc.)
    save_path = None
    if " d:" in command_lower or " d drive" in command_lower:
        save_path = "D:\\"
    elif " save to " in command_lower:
        import re
        match = re.search(r"save to\s+([a-z]:\\)", command_lower)
        if match:
            save_path = match.group(1)
    
    # Clean up the topic
    if topic:
        # Remove save location info from topic
        topic = topic.replace(" save to d:", "").replace(" save to d drive", "").strip()
        # Remove extra words and clean up
        topic = topic.replace(" and save", "").replace(" and create", "").strip()
        # Take only the first reasonable length (probably a topic name)
        topic = " ".join(topic.split()[:5])  # Limit to 5 words for topic
    
    return topic, save_path

def get_function_for_command(command: str):
    """Get the appropriate function for a natural language command"""
    command_lower = command.lower().strip()

    # Check for PowerPoint commands first
    for key, func_name in FUNCTION_MAPPINGS.items():
        if isinstance(func_name, str) and func_name == "create_powerpoint_presentation_function":
            if key in command_lower:
                # Extract topic and save path for PowerPoint creation
                topic, save_path = _parse_powerpoint_command(command)
                if topic:
                    return lambda: create_powerpoint_presentation(topic, save_path=save_path)
                break

    # Check for write commands
    for key, func_name in FUNCTION_MAPPINGS.items():
        if isinstance(func_name, str) and func_name == "write_text_function":
            if key in command_lower:
                # Extract parameters for writing text
                content, filename, save_path = _parse_write_command(command)
                if content:
                    return lambda: write_text_to_file(content, filename, save_path)
                break

    # Check for scraping commands first
    for key, func_name in FUNCTION_MAPPINGS.items():
        if isinstance(func_name, str) and func_name == "scrape_info_about_function":
            if key in command_lower:
                # Extract parameters for scraping
                search_term, info_type = _parse_scraping_command(command)
                if search_term:
                    return lambda: scrape_info_about(search_term, info_type)
                break

    # Direct mapping
    if command_lower in FUNCTION_MAPPINGS:
        func = FUNCTION_MAPPINGS[command_lower]
        if isinstance(func, str):
            return None  # Skip string references in direct mapping
        return func

    # Partial matching
    for key, func in FUNCTION_MAPPINGS.items():
        if isinstance(func, str):  # Skip string references
            continue
        if key in command_lower:
            return func

    return None


def list_available_functions() -> dict:
    """List all available functions with descriptions"""
    functions = {}

    # Get all functions from this module
    import inspect
    current_module = inspect.getmodule(inspect.currentframe())

    for name, obj in inspect.getmembers(current_module):
        if inspect.isfunction(obj) and not name.startswith('_'):
            doc = inspect.getdoc(obj) or "No description available"
            functions[name] = {
                'function': obj,
                'description': doc.split('\n')[0],  # First line of docstring
                'parameters': list(inspect.signature(obj).parameters.keys())
            }

    return functions

# ========================
# WEB AND MEDIA FUNCTIONS (Migrated from FunctionExecutor)
# ========================

def open_website(url: str) -> bool:
    """Opens a website in the default browser."""
    import webbrowser
    try:
        if not url.startswith("http"):
            url = "https://" + url
        webbrowser.open(url)
        print(f"Opening website: {url}")
        return True
    except Exception as e:
        print(f"Error opening website: {e}")
        return False

def google_search(query: str) -> bool:
    """Performs a Google search."""
    import webbrowser
    import urllib.parse
    try:
        if query:
            url = f"https://www.google.com/search?q={urllib.parse.quote(query)}"
        else:
            url = "https://www.google.com"
        webbrowser.open(url)
        print(f"Searching Google for: {query}")
        return True
    except Exception as e:
        print(f"Error searching Google: {e}")
        return False

def play_youtube(query: str = "") -> bool:
    """Opens YouTube, optionally searching for a query."""
    import webbrowser
    import urllib.parse
    try:
        if query:
            url = f"https://www.youtube.com/results?search_query={urllib.parse.quote(query)}"
        else:
            url = "https://www.youtube.com"
        webbrowser.open(url)
        print(f"Opening YouTube: {query if query else 'Home'}")
        return True
    except Exception as e:
        print(f"Error opening YouTube: {e}")
        return False

def play_spotify(query: str = "") -> bool:
    """Opens Spotify (App or Web), optionally searching."""
    import webbrowser
    import urllib.parse
    import subprocess
    try:
        # Try app first
        try:
            subprocess.Popen("spotify", shell=True)
            if query:
                # Also open web search as backup/complement
                url = f"https://open.spotify.com/search/{urllib.parse.quote(query)}"
                webbrowser.open(url)
            print(f"Opening Spotify: {query if query else 'Home'}")
            return True
        except:
            # Web fallback
            url = f"https://open.spotify.com/search/{urllib.parse.quote(query)}" if query else "https://open.spotify.com"
            webbrowser.open(url)
            print(f"Opening Spotify Web: {query}")
            return True
    except Exception as e:
        print(f"Error opening Spotify: {e}")
        return False

def media_control(action: str) -> bool:
    """Controls media playback (play_pause, next, previous)."""
    import ctypes
    try:
        # Virtual Key Codes
        keys = {
            "play_pause": 0xB3,
            "next": 0xB0,
            "previous": 0xB1
        }
        
        vk_code = keys.get(action)
        if not vk_code:
            print(f"Unknown media action: {action}")
            return False
            
        ctypes.windll.user32.keybd_event(vk_code, 0, 0, 0)
        ctypes.windll.user32.keybd_event(vk_code, 0, 2, 0)
        print(f"Media action: {action}")
        return True
    except Exception as e:
        print(f"Error checking media control: {e}")
        return False

# ========================
# UTILITY FUNCTIONS
# ========================

def create_file(file_name: str, content: str = "", location: str = "") -> bool:
    """Creates a file with content at a specified location (default: Desktop)."""
    import os
    try:
        # Determine base path
        if location:
            loc = location.lower().replace(" drive", "").replace("drive ", "").strip()
            if len(loc) == 1 and loc.isalpha():
                base_path = f"{loc.upper()}:\\"
            elif loc.endswith(":"):
                base_path = loc.upper() + "\\"
            else:
                base_path = location
        else:
            base_path = os.path.join(os.path.expanduser("~"), "Desktop")
        
        # Ensure valid path
        if not os.path.exists(base_path):
             os.makedirs(base_path, exist_ok=True)

        # Ensure extension
        if "." not in file_name:
            file_name += ".txt"
            
        file_path = os.path.join(base_path, file_name)
        
        with open(file_path, 'w', encoding='utf-8') as f:
            f.write(content)
            
        print(f"Created file: {file_path}")
        return True
    except Exception as e:
        print(f"Error creating file: {e}")
        return False

def open_calculator() -> bool:
    """Opens the Windows Calculator."""
    try:
        subprocess.Popen("calc", shell=True)
        return True
    except Exception as e:
        print(f"Error opening calculator: {e}")
        return False

def get_current_date_time(format_str: str = "both") -> str:
    """Returns current date, time, or both."""
    from datetime import datetime
    now = datetime.now()
    if format_str == "time":
        return now.strftime("%I:%M %p")
    elif format_str == "date":
        return now.strftime("%A, %B %d, %Y")
    else:
        return now.strftime("%A, %B %d, %Y %I:%M %p")

# ========================
# APP CONTROL (Migrated from FunctionExecutor)
# ========================

def open_application(app_name: str) -> bool:
    """Opens an application by name with support for Store apps and URLs."""
    import os
    import subprocess
    import webbrowser
    
    app_name = app_name.lower().strip()
    
    # Map common app names to executables/commands
    # Format: app_name -> (executable, is_store_app)
    app_map = {
        "chrome": ("chrome", False),
        "google chrome": ("chrome", False),
        "firefox": ("firefox", False),
        "edge": ("msedge", False),
        "microsoft edge": ("msedge", False),
        "notepad": ("notepad", False),
        "calculator": ("calc", False),
        "calc": ("calc", False),
        "paint": ("mspaint", False),
        "word": ("winword", False),
        "excel": ("excel", False),
        "powerpoint": ("powerpnt", False),
        "cmd": ("cmd", False),
        "command prompt": ("cmd", False),
        "terminal": ("wt", False),
        "powershell": ("powershell", False),
        "settings": ("ms-settings:", True),
        "control panel": ("control", False),
        "task manager": ("taskmgr", False),
        "spotify": ("spotify", True),
        "vscode": ("code", False),
        "vs code": ("code", False),
        "visual studio code": ("code", False),
        "whatsapp": ("whatsapp", True),
        "telegram": ("telegram", True),
        "discord": ("discord", False),
        "slack": ("slack", False),
        "vlc": ("vlc", False),
        "zoom": ("zoom", False),
        "yt": ("https://www.youtube.com", False),
        "youtube": ("https://www.youtube.com", False),
    }
    
    app_info = app_map.get(app_name, (app_name, False))
    executable, is_store_app = app_info
    
    try:
        if executable.startswith("http"):
             webbrowser.open(executable)
             print(f"Opened URL: {executable}")
             return True

        if executable.startswith("ms-"):
            # Windows URI scheme
            os.system(f'start "" "{executable}"')
            print(f"Launched {app_name}")
            return True
        
        if is_store_app:
            # Try multiple methods for store apps
            # Method 1: Use start command with app name
            result = os.system(f'start "" "{executable}:"')
            if result == 0:
                print(f"Launched {app_name}")
                return True
            
            # Method 2: Try explorer shell:AppsFolder (Specific for Spotify)
            if executable == "spotify":
                 subprocess.run(
                    f'explorer.exe shell:AppsFolder\\SpotifyAB.SpotifyMusic_zpdnekdrzrea0!Spotify',
                    shell=True
                )
                 
            # Method 3: Web fallback for Spotify
            if executable == "spotify":
                webbrowser.open("https://open.spotify.com")
                print(f"Opened Spotify Web Player")
                return True
            
            print(f"Launched {app_name}")
            return True
        
        # Regular apps - try start command first
        result = os.system(f'start "" "{executable}"')
        if result == 0:
            print(f"Launched {app_name}")
            return True
        
        # Fallback: direct execution
        subprocess.Popen(executable, shell=True)
        print(f"Launched {app_name}")
        return True
        
    except Exception as e:
        print(f"Failed to open {app_name}: {e}")
        return False

def close_application(app_name: str) -> bool:
    """Closes an application by name."""
    import subprocess
    app_name = app_name.lower().strip()
    try:
        subprocess.run(f'taskkill /IM {app_name}.exe /F', shell=True, check=True, stdout=subprocess.DEVNULL, stderr=subprocess.DEVNULL)
        print(f"Closed {app_name}")
        return True
    except:
        try:
            subprocess.run(f'taskkill /FI "WINDOWTITLE eq {app_name}*" /F', shell=True, stdout=subprocess.DEVNULL, stderr=subprocess.DEVNULL)
            print(f"Closed {app_name} via window title")
            return True
        except Exception as e:
            print(f"Failed to close {app_name}: {e}")
            return False

def adjust_system_volume(change: int) -> bool:
    """Adjusts system volume by a relative amount."""
    try:
        current = get_current_volume()
        new_level = max(0, min(100, current + change))
        return set_system_volume(new_level)
    except Exception as e:
        print(f"Error adjusting volume: {e}")
        return False

def get_current_time() -> str:
    """
    Returns the current date and time formatted as a string.
    
    Returns:
        Formatted date time string "YYYY-MM-DD HH:MM:SS"
    """
    import datetime
    now = datetime.datetime.now()
    return now.strftime("%Y-%m-%d %H:%M:%S")
